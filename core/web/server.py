"""
NeuroClaw Web UI Server

Serves a browser-based chat interface at http://localhost:7080 by default.

Usage
-----
    # Preferred: via the main agent entry point
    python core/agent/main.py --web [--port 7080] [--host 127.0.0.1]

    # Or run the web server directly
    python core/web/server.py [--port 7080] [--host 127.0.0.1]

Dependencies (install via the setup wizard or manually)
---------------------------------------------------------
    pip install "fastapi[standard]" uvicorn
"""
from __future__ import annotations

import argparse
import asyncio
import importlib.util
import io
import json
import os
import subprocess
import time
import re
import queue as stdlib_queue
import sys
import threading
import urllib.error
import urllib.request
from pathlib import Path
from typing import Any

REPO_ROOT = Path(__file__).parent.parent.parent.resolve()
STATIC_DIR = Path(__file__).parent / "static"
SKILL_SUMMARIES_FILE = STATIC_DIR / "skill_summaries.json"
SHELL_STATUS_FILE = Path("/tmp/neuroclaw_claw_shell_status.json")
AGENT_SHELL_STATUS_FILE = Path("/tmp/neuroclaw_agent_shell_status.json")

DEFAULT_HOST = "0.0.0.0"
DEFAULT_PORT = 7080
ATTACHMENT_MAX_FILE_BYTES = 15 * 1024 * 1024
ATTACHMENT_MAX_EMBED_CHARS = 12000
SUPPORTED_ATTACHMENT_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".json", ".yaml", ".yml", ".csv", ".tsv",
    ".py", ".js", ".ts", ".tsx", ".jsx", ".sh", ".bash", ".zsh", ".sql",
    ".html", ".css", ".xml", ".log", ".rst", ".ini", ".toml", ".cfg",
    ".pdf", ".docx", ".xlsx", ".pptx",
}


def _fallback_title_from_user_text(text: str) -> str:
    """Create a short deterministic title when LLM title generation fails."""
    cleaned = re.sub(r"\s+", " ", str(text or "")).strip()
    if not cleaned:
        return "New Chat"
    words = cleaned.split(" ")[:8]
    title = " ".join(words)
    return title[:64] if title else "New Chat"


def _sanitize_title(raw: str, user_text: str) -> str:
    """Normalize model output into a single short title line."""
    t = str(raw or "").strip()
    if not t:
        return _fallback_title_from_user_text(user_text)

    t = t.splitlines()[0].strip()
    t = t.strip("\"'` ")
    t = re.sub(r"^[\-\*\d\.)\s]+", "", t)
    t = re.sub(r"\s+", " ", t).strip(" .:-")
    if not t:
        return _fallback_title_from_user_text(user_text)
    return t[:64]


def _summarize_web_tool_events(events: list[dict[str, Any]]) -> list[dict[str, Any]]:
    """Return compact tool events that are safe to show in the web timeline."""
    compact: list[dict[str, Any]] = []
    for idx, event in enumerate(events, start=1):
        if not isinstance(event, dict):
            continue
        result = event.get("result") if isinstance(event.get("result"), dict) else {}
        stdout = str(result.get("stdout", "") or result.get("output", "") or "")
        stderr = str(result.get("stderr", "") or result.get("error", "") or "")
        compact.append(
            {
                "id": idx,
                "tool": str(event.get("tool", "tool")),
                "command": str(event.get("command", "")),
                "executed": bool(event.get("executed", False)),
                "success": bool(event.get("success", False)),
                "stdout_preview": stdout[:1200],
                "stderr_preview": stderr[:1200],
                "skills_used": event.get("skills_used", []) if isinstance(event.get("skills_used"), list) else [],
            }
        )
    return compact


def _workspace_change_summary() -> list[dict[str, str]]:
    """Summarize current git working tree changes for the web timeline."""
    try:
        proc = subprocess.run(
            ["git", "status", "--short"],
            cwd=REPO_ROOT,
            capture_output=True,
            text=True,
            encoding="utf-8",
            errors="replace",
            timeout=5,
        )
    except Exception:
        return []
    if proc.returncode != 0:
        return []
    changes: list[dict[str, str]] = []
    for line in proc.stdout.splitlines():
        if not line.strip():
            continue
        status = line[:2].strip() or "modified"
        path = line[3:].strip() if len(line) > 3 else line.strip()
        changes.append({"status": status, "path": path})
    return changes[:40]


def _safe_skill_summary_fallback(skill_name: str, description: str) -> dict[str, str]:
    """Return conservative bilingual fallback summary when LLM summarization fails."""
    base_en = (description or "").strip()
    if not base_en:
        base_en = f"{skill_name} provides a specialized workflow for task execution in NeuroClaw."
    base_en = re.sub(r"\s+", " ", base_en).strip()
    if len(base_en) > 220:
        base_en = base_en[:217].rstrip() + "..."
    base_zh = f"该技能围绕「{skill_name}」提供专用流程支持，可用于相关任务的执行与组织。"
    return {"en": base_en, "zh": base_zh}


def _parse_bilingual_summary_json(raw: str) -> dict[str, str] | None:
    """Parse model output into {'en':..., 'zh':...} with defensive handling."""
    text = str(raw or "").strip()
    if not text:
        return None

    try:
        obj = json.loads(text)
        en = str(obj.get("en", "")).strip()
        zh = str(obj.get("zh", "")).strip()
        if en and zh:
            return {"en": en, "zh": zh}
    except Exception:
        pass

    match = re.search(r"\{[\s\S]*\}", text)
    if match:
        try:
            obj = json.loads(match.group(0))
            en = str(obj.get("en", "")).strip()
            zh = str(obj.get("zh", "")).strip()
            if en and zh:
                return {"en": en, "zh": zh}
        except Exception:
            return None
    return None


def _strip_frontmatter(text: str) -> str:
    """Remove optional YAML front-matter from SKILL.md text."""
    return re.sub(r"^---\s*\n.*?\n---\s*\n", "", text, flags=re.DOTALL)


def _read_shell_status() -> dict[str, Any] | None:
    """Read transient shell status from agent-shell first, then tmux claw-shell."""
    if AGENT_SHELL_STATUS_FILE.exists():
        try:
            data = json.loads(AGENT_SHELL_STATUS_FILE.read_text(encoding="utf-8"))
            command = str(data.get("command", "")).strip()
            started_at = data.get("started_at")
            pid = int(data.get("pid", 0))
            if command and isinstance(started_at, (int, float)) and pid > 0:
                alive = False
                try:
                    os.kill(pid, 0)
                    alive = True
                except Exception:
                    alive = False

                if alive:
                    return {
                        "active": True,
                        "source": "agent_shell",
                        "command": command,
                        "current_command": "agent_shell",
                        "started_at": int(started_at),
                        "elapsed_ms": max(0, int((time.time() * 1000) - int(started_at))),
                    }
        except Exception:
            pass
        try:
            AGENT_SHELL_STATUS_FILE.unlink()
        except Exception:
            pass

    if not SHELL_STATUS_FILE.exists():
        return None

    try:
        data = json.loads(SHELL_STATUS_FILE.read_text(encoding="utf-8"))
    except Exception:
        try:
            SHELL_STATUS_FILE.unlink()
        except Exception:
            pass
        return None

    command = str(data.get("command", "")).strip()
    started_at = data.get("started_at")
    if not command or not isinstance(started_at, (int, float)):
        try:
            SHELL_STATUS_FILE.unlink()
        except Exception:
            pass
        return None

    try:
        proc = subprocess.run(
            ["tmux", "display-message", "-p", "-t", "claw", "#{pane_current_command}"],
            capture_output=True,
            text=True,
            check=False,
            timeout=2,
        )
        current = (proc.stdout or "").strip().lower()
    except Exception:
        current = ""

    shell_names = {"bash", "zsh", "fish", "sh", "dash", "tmux"}
    active = bool(current) and current not in shell_names
    elapsed_ms = max(0, int((time.time() * 1000) - int(started_at)))

    if not active:
        try:
            SHELL_STATUS_FILE.unlink()
        except Exception:
            pass
        return None

    return {
        "active": True,
        "command": command,
        "current_command": current,
        "started_at": int(started_at),
        "elapsed_ms": elapsed_ms,
    }


def _skill_md_excerpt(skill_md: Path, max_chars: int = 1800) -> str:
    """Return a compact excerpt from SKILL.md for prompt context injection."""
    try:
        raw = skill_md.read_text(encoding="utf-8")
    except Exception:
        return ""
    body = _strip_frontmatter(raw)
    body = re.sub(r"```[\s\S]*?```", "", body)
    body = re.sub(r"\n{3,}", "\n\n", body).strip()
    if len(body) > max_chars:
        body = body[:max_chars].rstrip() + "\n..."
    return body


def _selected_skills_context(selected_names: list[str], skills: list[dict[str, Any]]) -> str:
    """Build selected-skill references from discovered SKILL.md files."""
    if not selected_names:
        return ""

    by_name: dict[str, dict[str, Any]] = {}
    for s in skills:
        n = str(s.get("name", "")).strip().lower()
        if n:
            by_name[n] = s

    chunks: list[str] = []
    for raw_name in selected_names:
        key = raw_name.strip().lower()
        if not key:
            continue
        s = by_name.get(key)
        if not s:
            continue
        excerpt = _skill_md_excerpt(Path(s.get("skill_md")))
        if not excerpt:
            continue
        chunks.append(
            f"[Skill: {s.get('name', raw_name)}]\n"
            f"Description: {s.get('description', '')}\n"
            f"SKILL.md excerpt:\n{excerpt}"
        )

    return "\n\n".join(chunks)


def _attachment_extension(filename: str) -> str:
    return Path(str(filename or "")).suffix.lower().strip()


def _truncate_text(text: str, max_chars: int = ATTACHMENT_MAX_EMBED_CHARS) -> tuple[str, bool]:
    raw = str(text or "")
    if len(raw) <= max_chars:
        return raw, False
    return raw[:max_chars].rstrip() + "\n[Content truncated]", True


def _decode_text_attachment(raw: bytes) -> str:
    for encoding in ("utf-8", "utf-8-sig", "gb18030", "latin-1"):
        try:
            return raw.decode(encoding)
        except Exception:
            continue
    return raw.decode("utf-8", errors="replace")


def _extract_pdf_text(raw: bytes) -> str:
    from pypdf import PdfReader  # type: ignore

    reader = PdfReader(io.BytesIO(raw))
    chunks: list[str] = []
    for page in reader.pages[:120]:
        txt = (page.extract_text() or "").strip()
        if txt:
            chunks.append(txt)
    return "\n\n".join(chunks)


def _extract_docx_text(raw: bytes) -> str:
    from docx import Document  # type: ignore

    doc = Document(io.BytesIO(raw))
    chunks: list[str] = []

    for para in doc.paragraphs:
        txt = (para.text or "").strip()
        if txt:
            chunks.append(txt)

    for table in doc.tables:
        for row in table.rows:
            vals = [str(cell.text or "").strip() for cell in row.cells]
            vals = [v for v in vals if v]
            if vals:
                chunks.append(" | ".join(vals))

    return "\n".join(chunks)


def _extract_xlsx_text(raw: bytes) -> str:
    from openpyxl import load_workbook  # type: ignore

    wb = load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
    lines: list[str] = []
    for ws in wb.worksheets[:12]:
        lines.append(f"[Sheet] {ws.title}")
        row_count = 0
        for row in ws.iter_rows(max_row=200, max_col=30, values_only=True):
            cells = ["" if v is None else str(v).strip() for v in row]
            if not any(cells):
                continue
            lines.append("\t".join(cells))
            row_count += 1
            if row_count >= 200:
                lines.append("[Rows truncated]")
                break
        lines.append("")
    return "\n".join(lines).strip()


def _extract_pptx_text(raw: bytes) -> str:
    from pptx import Presentation  # type: ignore

    prs = Presentation(io.BytesIO(raw))
    lines: list[str] = []
    for idx, slide in enumerate(prs.slides[:80], start=1):
        lines.append(f"[Slide {idx}]")
        for shape in slide.shapes:
            txt = ""
            try:
                txt = (shape.text or "").strip()  # type: ignore[attr-defined]
            except Exception:
                txt = ""
            if txt:
                lines.append(txt)
        lines.append("")
    return "\n".join(lines).strip()


def _parse_attachment_content(filename: str, content_type: str, raw: bytes) -> dict[str, Any]:
    ext = _attachment_extension(filename)
    if ext not in SUPPORTED_ATTACHMENT_EXTENSIONS:
        return {
            "ok": False,
            "error": (
                f"Unsupported extension '{ext or 'unknown'}'. "
                "Supported: " + ", ".join(sorted(SUPPORTED_ATTACHMENT_EXTENSIONS))
            ),
        }

    if ext in {
        ".txt", ".md", ".markdown", ".json", ".yaml", ".yml", ".csv", ".tsv",
        ".py", ".js", ".ts", ".tsx", ".jsx", ".sh", ".bash", ".zsh", ".sql",
        ".html", ".css", ".xml", ".log", ".rst", ".ini", ".toml", ".cfg",
    }:
        text = _decode_text_attachment(raw)
    elif ext == ".pdf":
        try:
            text = _extract_pdf_text(raw)
        except ImportError:
            return {"ok": False, "error": "PDF parser missing. Install: pypdf"}
    elif ext == ".docx":
        try:
            text = _extract_docx_text(raw)
        except ImportError:
            return {"ok": False, "error": "DOCX parser missing. Install: python-docx"}
    elif ext == ".xlsx":
        try:
            text = _extract_xlsx_text(raw)
        except ImportError:
            return {"ok": False, "error": "XLSX parser missing. Install: openpyxl"}
    elif ext == ".pptx":
        try:
            text = _extract_pptx_text(raw)
        except ImportError:
            return {"ok": False, "error": "PPTX parser missing. Install: python-pptx"}
    else:
        return {"ok": False, "error": "Unsupported file type"}

    cleaned = str(text or "").strip()
    if not cleaned:
        return {"ok": False, "error": "No readable text found in file"}

    truncated_text, truncated = _truncate_text(cleaned)
    return {
        "ok": True,
        "text": truncated_text,
        "truncated": truncated,
        "ext": ext,
        "content_type": content_type,
    }


def _normalize_skill_token(text: str) -> str:
    """Normalize a skill token for robust mention matching."""
    t = str(text or "").lower()
    t = re.sub(r"[^a-z0-9]+", "", t)
    return t


def _infer_skills_from_user_text(user_text: str, skills: list[dict[str, Any]]) -> list[str]:
    """
    Infer referenced skills from free-form user text.

    Matches both skill name and skill directory name in normalized form.
    """
    norm_msg = _normalize_skill_token(user_text)
    if not norm_msg:
        return []

    inferred: list[str] = []
    for s in skills:
        skill_name = str(s.get("name", "")).strip()
        if not skill_name:
            continue
        dir_name = ""
        try:
            dir_name = Path(s.get("path")).name
        except Exception:
            dir_name = ""

        candidates = [skill_name, dir_name]
        for c in candidates:
            token = _normalize_skill_token(c)
            if token and token in norm_msg:
                inferred.append(skill_name)
                break

    # Preserve order and remove duplicates
    seen: set[str] = set()
    out: list[str] = []
    for name in inferred:
        key = name.lower()
        if key in seen:
            continue
        seen.add(key)
        out.append(name)
    return out[:5]


# ── Module import helpers ──────────────────────────────────────────────────────

def _import_from_path(module_name: str, path: Path) -> Any:
    """Import a Python module from an absolute file path (handles hyphenated dirs)."""
    spec = importlib.util.spec_from_file_location(module_name, path)
    if spec is None or spec.loader is None:
        raise ImportError(f"Cannot load module from {path}")
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)  # type: ignore[union-attr]
    return mod


# ── Dependency check ───────────────────────────────────────────────────────────

def _require_webdeps() -> None:
    """Raise a descriptive RuntimeError if fastapi/uvicorn are not installed."""
    missing = []
    for pkg in ("fastapi", "uvicorn"):
        try:
            __import__(pkg)
        except ImportError:
            missing.append(pkg)
    if missing:
        raise RuntimeError(
            f"Web UI dependencies not installed: {', '.join(missing)}\n"
            "Install with:  pip install 'fastapi[standard]' uvicorn\n"
            "Or re-run the installer:  python installer/setup.py"
        )


# ── Streaming helpers ──────────────────────────────────────────────────────────

async def _stream_openai(
    websocket: Any, llm_client: Any, model: str, history: list[dict]
) -> str:
    """
    Stream an OpenAI response chunk-by-chunk.

    Uses a producer thread + asyncio queue so the blocking OpenAI iterator
    does not stall the event loop while still delivering incremental updates
    to the browser.
    """
    q: stdlib_queue.Queue[tuple[str, str | None]] = stdlib_queue.Queue()

    def _produce() -> None:
        try:
            stream = llm_client.chat.completions.create(
                model=model, messages=history, stream=True
            )
            for chunk in stream:
                content = ""
                if chunk.choices and chunk.choices[0].delta:
                    content = chunk.choices[0].delta.content or ""
                if content:
                    q.put(("chunk", content))
        except Exception as exc:
            q.put(("error", str(exc)))
        finally:
            q.put(("done", None))

    threading.Thread(target=_produce, daemon=True).start()

    full = ""
    while True:
        try:
            kind, data = q.get_nowait()
        except stdlib_queue.Empty:
            await asyncio.sleep(0.01)
            continue

        if kind == "chunk":
            full += data  # type: ignore[operator]
            await websocket.send_text(json.dumps({"type": "chunk", "content": data}))
        elif kind == "error":
            raise RuntimeError(data)
        else:  # "done"
            break

    return full


async def _respond(websocket: Any, session: Any) -> str:
    """
    Generate a reply for the latest message in session.history.

    Streams chunks to the browser for OpenAI-compatible providers; falls back
    to asyncio.to_thread for other backends (Anthropic, local).  Always sends a final
    ``{"type": "done", "content": "…"}`` frame.
    """
    provider = session.env.get("llm_backend", {}).get("provider", "openai")
    model = session.env.get("llm_backend", {}).get("model", "gpt-4o")
    from core.llm.provider_profiles import is_openai_compatible_provider

    if is_openai_compatible_provider(provider) and session._llm is not None:
        full = await _stream_openai(websocket, session._llm, model, session.history)
    else:
        # Non-streaming fallback: run blocking _chat() in a thread pool
        full = await asyncio.to_thread(session._chat)

    await websocket.send_text(json.dumps({"type": "done", "content": full}))
    return full


# ── FastAPI application factory ────────────────────────────────────────────────

def create_app() -> Any:
    """Build and return the FastAPI application object."""
    _require_webdeps()

    from fastapi import Body, FastAPI, File, UploadFile, WebSocket, WebSocketDisconnect  # type: ignore
    from fastapi.responses import FileResponse, HTMLResponse, JSONResponse  # type: ignore
    from fastapi.staticfiles import StaticFiles  # type: ignore

    # Ensure repo root is on sys.path so `from core.agent.main import …` resolves
    if str(REPO_ROOT) not in sys.path:
        sys.path.insert(0, str(REPO_ROOT))

    from core.agent.main import (  # type: ignore[import]
        AgentSession,
        build_llm_client,
        load_environment,
        save_environment,
    )

    # SkillLoader lives in core/skill_loader/, so we use importlib for dynamic loading.
    _loader_mod = _import_from_path(
        "neuroclaw_skill_loader",
        REPO_ROOT / "core" / "skill_loader" / "loader.py",
    )
    SkillLoader = _loader_mod.SkillLoader

    def _load_offline_skill_summaries() -> dict[str, dict[str, str]]:
        try:
            raw = json.loads(SKILL_SUMMARIES_FILE.read_text(encoding="utf-8"))
        except Exception:
            return {}
        if not isinstance(raw, dict):
            return {}
        out: dict[str, dict[str, str]] = {}
        for name, summary in raw.items():
            if not isinstance(name, str) or not isinstance(summary, dict):
                continue
            en = str(summary.get("en", "")).strip()
            zh = str(summary.get("zh", "")).strip()
            if en or zh:
                out[name.lower()] = {"en": en, "zh": zh}
        return out

    def _summary_for_skill(skill: dict[str, Any], offline: dict[str, dict[str, str]]) -> dict[str, str]:
        name = str(skill.get("name", "")).strip()
        summary = offline.get(name.lower()) if name else None
        if summary:
            return {
                "en": summary.get("en") or str(skill.get("summary_en") or skill.get("description") or ""),
                "zh": summary.get("zh") or str(skill.get("summary_zh") or ""),
            }
        return _safe_skill_summary_fallback(name, str(skill.get("description", "")))

    def _dedupe_model_catalog(models: list[dict[str, Any]]) -> list[dict[str, Any]]:
        out: list[dict[str, Any]] = []
        seen: set[tuple[str, str]] = set()
        for item in models:
            if not isinstance(item, dict):
                continue
            provider = str(item.get("provider") or "openai").strip() or "openai"
            model = str(item.get("model") or item.get("id") or item.get("name") or "").strip()
            if not model:
                continue
            key = (provider.lower(), model)
            if key in seen:
                continue
            seen.add(key)
            entry = dict(item)
            entry["provider"] = provider
            entry["model"] = model
            entry.setdefault("label", model)
            out.append(entry)
        return out

    def _fetch_openai_compatible_models(llm: dict[str, Any]) -> list[dict[str, Any]]:
        base_url = str(llm.get("base_url") or llm.get("baseUrl") or "").strip()
        if not base_url:
            return []
        api_key = ""
        api_key_env = str(llm.get("api_key_env") or "").strip()
        if api_key_env:
            api_key = os.environ.get(api_key_env, "")
        api_key = api_key or str(llm.get("api_key") or llm.get("apiKey") or "").strip()
        url = base_url.rstrip("/") + "/models"
        headers = {"Accept": "application/json"}
        if api_key:
            headers["Authorization"] = f"Bearer {api_key}"
        req = urllib.request.Request(url, headers=headers, method="GET")
        try:
            with urllib.request.urlopen(req, timeout=3) as resp:
                payload = json.loads(resp.read().decode("utf-8"))
        except (OSError, urllib.error.URLError, json.JSONDecodeError):
            return []
        data = payload.get("data") if isinstance(payload, dict) else None
        if not isinstance(data, list):
            return []
        provider = str(llm.get("provider") or "openai").strip() or "openai"
        models: list[dict[str, Any]] = []
        inherited_keys = (
            "base_url", "baseUrl", "api_key_env", "api_key", "apiKey",
            "default_headers", "headers", "tool_calling",
        )
        for raw in data:
            if not isinstance(raw, dict):
                continue
            model = str(raw.get("id") or raw.get("model") or raw.get("name") or "").strip()
            if not model:
                continue
            entry: dict[str, Any] = {"provider": provider, "model": model, "label": model}
            for key in inherited_keys:
                if key in llm:
                    entry[key] = llm[key]
            models.append(entry)
        return models

    def _runtime_available_models(llm: dict[str, Any]) -> list[dict[str, Any]]:
        configured = llm.get("available_models", [])
        base = configured if isinstance(configured, list) else []
        return _dedupe_model_catalog([*base, *_fetch_openai_compatible_models(llm)])

    app = FastAPI(title="NeuroClaw Web UI", docs_url=None, redoc_url=None)

    # ── Static files ────────────────────────────────────────────────────────────
    if STATIC_DIR.exists():
        app.mount("/static", StaticFiles(directory=str(STATIC_DIR)), name="static")
    materials_dir = REPO_ROOT / "materials"
    if materials_dir.exists():
        app.mount("/materials", StaticFiles(directory=str(materials_dir)), name="materials")

    # ── HTTP endpoints ──────────────────────────────────────────────────────────

    @app.get("/")
    async def root() -> Any:
        index = STATIC_DIR / "index.html"
        if index.exists():
            return FileResponse(str(index))
        return HTMLResponse(
            "<h1>NeuroClaw Web UI</h1>"
            f"<p>Static files not found at <code>{STATIC_DIR}</code>.</p>",
            status_code=500,
        )

    @app.get("/api/health")
    async def health() -> dict:
        return {"status": "ok"}

    @app.get("/api/shell/status")
    async def shell_status() -> dict:
        status = _read_shell_status()
        if not status:
            return {"active": False}
        return status

    @app.get("/api/skills")
    async def list_skills() -> dict:
        loader = SkillLoader(REPO_ROOT / "skills")
        skills = loader.load_all()
        offline_summaries = _load_offline_skill_summaries()
        return {
            "skills": [
                {
                    "name": s["name"],
                    "description": s.get("description", ""),
                    "summary_en": _summary_for_skill(s, offline_summaries)["en"],
                    "summary_zh": _summary_for_skill(s, offline_summaries)["zh"],
                    "layer": s.get("layer", ""),
                    "skill_type": s.get("skill_type", ""),
                    "dependencies": s.get("dependencies", []),
                    "complementary_skills": s.get("complementary_skills", []),
                }
                for s in skills
            ]
        }

    @app.post("/api/attachments/parse")
    async def parse_attachments(files: list[UploadFile] = File(...)) -> Any:
        if not files:
            return JSONResponse({"type": "error", "message": "No files uploaded"}, status_code=400)

        parsed_files: list[dict[str, Any]] = []
        for upload in files[:24]:
            name = str(upload.filename or "untitled")
            content_type = str(upload.content_type or "unknown")
            raw = await upload.read()
            size = len(raw)

            item: dict[str, Any] = {
                "name": name,
                "size": size,
                "content_type": content_type,
            }

            if size > ATTACHMENT_MAX_FILE_BYTES:
                item.update(
                    {
                        "ok": False,
                        "error": (
                            f"File too large ({size} bytes). "
                            f"Limit is {ATTACHMENT_MAX_FILE_BYTES} bytes."
                        ),
                    }
                )
                parsed_files.append(item)
                continue

            try:
                item.update(_parse_attachment_content(name, content_type, raw))
            except Exception as exc:
                item.update({"ok": False, "error": f"Parse failed: {exc}"})
            parsed_files.append(item)

        return {
            "type": "done",
            "files": parsed_files,
            "max_file_bytes": ATTACHMENT_MAX_FILE_BYTES,
            "supported_extensions": sorted(SUPPORTED_ATTACHMENT_EXTENSIONS),
        }

    @app.get("/api/env")
    async def get_env() -> dict:
        """Return non-sensitive parts of the runtime environment config."""
        env = load_environment()
        llm = env.get("llm_backend", {})
        provider = llm.get("provider", "unknown")
        api_key_env = llm.get("api_key_env", "")
        direct_api_key = llm.get("api_key") or llm.get("apiKey")
        api_key_present = bool(
            llm.get("no_api_key_required")
            or (isinstance(direct_api_key, str) and direct_api_key.strip())
            or (api_key_env and __import__('os').environ.get(api_key_env))
        )
        available_models = _runtime_available_models(llm if isinstance(llm, dict) else {})
        return {
            "provider": provider,
            "model": llm.get("model", "unknown"),
            "base_url": llm.get("base_url") or llm.get("baseUrl") or llm.get("local_endpoint") or "",
            "available_models": available_models,
            "cuda_device": env.get("cuda", {}).get("device", "cpu"),
            "setup_type": env.get("setup_type", "unknown"),
            "conda_env": env.get("conda_env"),
            "api_key_present": api_key_present,
        }

    @app.post("/api/env/model")
    async def set_model(payload: dict) -> Any:
        """Switch current provider/model to one of the configured options."""
        provider = str(payload.get("provider", "")).strip()
        model = str(payload.get("model", "")).strip()
        if not provider or not model:
            return JSONResponse(
                {"type": "error", "message": "provider and model are required"},
                status_code=400,
            )

        env = load_environment()
        llm = env.setdefault("llm_backend", {})
        available_models = _runtime_available_models(llm)
        selected_model = next(
            (
                item
                for item in available_models
                if isinstance(item, dict)
                and str(item.get("provider", "")).strip() == provider
                and str(item.get("model", item.get("id", item.get("name", "")))).strip() == model
            ),
            None,
        )
        if selected_model is None:
            return JSONResponse(
                {"type": "error", "message": "Requested provider/model is not configured"},
                status_code=400,
            )

        llm["provider"] = provider
        llm["model"] = model
        for key in (
            "base_url",
            "baseUrl",
            "api_key_env",
            "api_key",
            "apiKey",
            "local_endpoint",
            "openai_compatible",
            "no_api_key_required",
            "default_headers",
            "headers",
            "tool_calling",
        ):
            if key in selected_model:
                llm[key] = selected_model[key]
        save_environment(env)
        return {
            "type": "done",
            "provider": provider,
            "model": model,
            "available_models": _runtime_available_models(llm),
        }

    @app.post("/api/chat")
    async def chat_http(payload: dict) -> Any:
        """HTTP fallback for chat when WebSocket is unavailable."""
        user_text = str(payload.get("message", "")).strip()
        raw_history = payload.get("history", [])
        raw_selected_skills = payload.get("selected_skills", [])
        if not user_text:
            return JSONResponse({"type": "error", "message": "Empty message"}, status_code=400)

        session = AgentSession()

        try:
            loader = SkillLoader(REPO_ROOT / "skills")
            skills = loader.load_all()
        except Exception:
            skills = []

        selected_skills: list[str] = []
        if isinstance(raw_selected_skills, list):
            selected_skills = [
                str(x).strip() for x in raw_selected_skills if str(x).strip()
            ]
        if not selected_skills:
            selected_skills = _infer_skills_from_user_text(user_text, skills)

        env_file = REPO_ROOT / "neuroclaw_environment.json"
        if not env_file.exists():
            return JSONResponse({
                "type": "error",
                "message": "neuroclaw_environment.json not found. Run python installer/setup.py to configure NeuroClaw.",
            }, status_code=400)

        try:
            session.set_llm_client(build_llm_client(session.env))
        except Exception as exc:
            return JSONResponse({"type": "error", "message": f"LLM backend error: {exc}"}, status_code=500)

        soul_path = REPO_ROOT / "SOUL.md"
        soul = soul_path.read_text(encoding="utf-8") if soul_path.exists() else ""
        skill_names = ", ".join(s["name"] for s in skills)
        session.history = [{"role": "system", "content": f"{soul}\n\nLoaded skills: {skill_names}"}]

        if isinstance(raw_history, list):
            for msg in raw_history:
                if not isinstance(msg, dict):
                    continue
                role = str(msg.get("role", "")).strip().lower()
                content = str(msg.get("content", "")).strip()
                if role not in {"user", "assistant"} or not content:
                    continue
                session.history.append({"role": role, "content": content})

        selected_ctx = _selected_skills_context(selected_skills, skills)
        user_payload = user_text
        if selected_ctx:
            user_payload = (
                f"{user_text}\n\n"
                "[Selected skill references from local SKILL.md files]\n"
                f"{selected_ctx}"
            )

        session.history.append({"role": "user", "content": user_payload})

        reply = await asyncio.to_thread(session._chat)
        llm_cfg = session.env.get("llm_backend", {})
        provider_used = str(llm_cfg.get("provider", "unknown"))
        model_used = str(llm_cfg.get("model", "unknown"))
        return {
            "type": "done",
            "content": reply,
            "provider_used": provider_used,
            "model_used": model_used,
            "tool_events": _summarize_web_tool_events(getattr(session, "_tool_events", [])),
            "workspace_changes": _workspace_change_summary(),
        }

    @app.post("/api/chat/title")
    async def chat_title(payload: dict) -> Any:
        """Generate a concise chat title using the configured default model."""
        user_text = str(payload.get("user", "")).strip()
        assistant_text = str(payload.get("assistant", "")).strip()
        if not user_text and not assistant_text:
            return JSONResponse({"type": "error", "message": "Missing conversation content"}, status_code=400)

        env_file = REPO_ROOT / "neuroclaw_environment.json"
        if not env_file.exists():
            return JSONResponse({"type": "error", "message": "Environment not configured"}, status_code=400)

        session = AgentSession()
        try:
            session.set_llm_client(build_llm_client(session.env))
        except Exception as exc:
            return JSONResponse({"type": "error", "message": f"LLM backend error: {exc}"}, status_code=500)

        title_system = (
            "You are a conversation title generator. "
            "Return exactly one short title in plain text, 3-10 words, no quotes, no markdown."
        )
        convo = f"User message:\n{user_text}\n\nAssistant reply:\n{assistant_text}"
        session.history = [
            {"role": "system", "content": title_system},
            {"role": "user", "content": convo},
        ]

        try:
            raw_title = await asyncio.to_thread(session._chat)
            title = _sanitize_title(raw_title, user_text)
        except Exception:
            title = _fallback_title_from_user_text(user_text)

        return {"type": "done", "title": title}

    @app.post("/api/skills/summary")
    async def skill_summary(payload: dict) -> Any:
        """Summarize one SKILL.md into bilingual 1-5 sentence summaries."""
        skill_name = str(payload.get("name", "")).strip()
        if not skill_name:
            return JSONResponse({"type": "error", "message": "Missing skill name"}, status_code=400)

        try:
            loader = SkillLoader(REPO_ROOT / "skills")
            skills = loader.load_all()
        except Exception:
            skills = []

        target = None
        for s in skills:
            if str(s.get("name", "")).strip().lower() == skill_name.lower():
                target = s
                break
        if target is None:
            return JSONResponse({"type": "error", "message": f"Skill not found: {skill_name}"}, status_code=404)

        summary = _summary_for_skill(target, _load_offline_skill_summaries())
        return {"type": "done", "summary_en": summary["en"], "summary_zh": summary["zh"]}

    # ── Checkpoint API endpoints ─────────────────────────────────────────────

    @app.get("/api/checkpoints")
    async def list_checkpoints() -> Any:
        from core.checkpoint.manager import ShadowCheckpointManager
        mgr = ShadowCheckpointManager(repo_root=REPO_ROOT)
        cps = mgr.list_checkpoints(REPO_ROOT)
        return {"checkpoints": cps}

    @app.get("/api/checkpoints/{checkpoint_id}/diff")
    async def checkpoint_diff(checkpoint_id: str) -> Any:
        from core.checkpoint.manager import ShadowCheckpointManager
        mgr = ShadowCheckpointManager(repo_root=REPO_ROOT)
        try:
            diff = mgr.diff_checkpoint(REPO_ROOT, checkpoint_id)
        except Exception as exc:
            return JSONResponse({"error": str(exc)}, status_code=400)
        return diff

    @app.post("/api/checkpoints/{checkpoint_id}/restore")
    async def restore_checkpoint(checkpoint_id: str, payload: dict = {}) -> Any:
        from core.checkpoint.manager import ShadowCheckpointManager
        mgr = ShadowCheckpointManager(repo_root=REPO_ROOT)
        filepath = payload.get("filepath") if payload else None
        try:
            result = mgr.restore_checkpoint(REPO_ROOT, checkpoint_id, filepath=filepath)
        except Exception as exc:
            return JSONResponse({"error": str(exc)}, status_code=400)
        return {"type": "done", **result}

    @app.get("/api/checkpoints/{checkpoint_id}/files")
    async def checkpoint_files(checkpoint_id: str) -> Any:
        from core.checkpoint.manager import ShadowCheckpointManager
        mgr = ShadowCheckpointManager(repo_root=REPO_ROOT)
        try:
            files = mgr.get_files_at_checkpoint(REPO_ROOT, checkpoint_id)
        except Exception as exc:
            return JSONResponse({"error": str(exc)}, status_code=400)
        return {"files": files}

    # ── WebSocket chat endpoint ─────────────────────────────────────────────────

    @app.websocket("/ws/chat")
    async def chat_endpoint(websocket: WebSocket) -> None:
        try:
            await websocket.accept()
            print("[WS] Client connected", flush=True)

            # Create a per-connection agent session
            session = AgentSession()

            # Load skills
            try:
                loader = SkillLoader(REPO_ROOT / "skills")
                skills = loader.load_all()
            except Exception:
                skills = []

            # Send init metadata
            llm_cfg = session.env.get("llm_backend", {})
            await websocket.send_text(json.dumps({
                "type": "init",
                "skills": [
                    {
                        "name": s["name"],
                        "description": s.get("description", ""),
                        "summary_en": s.get("summary_en", ""),
                        "summary_zh": s.get("summary_zh", ""),
                        "layer": s.get("layer", ""),
                        "skill_type": s.get("skill_type", ""),
                        "dependencies": s.get("dependencies", []),
                        "complementary_skills": s.get("complementary_skills", []),
                    }
                    for s in skills
                ],
                "provider": llm_cfg.get("provider", "unconfigured"),
                "model": llm_cfg.get("model", "unconfigured"),
            }))

            env_file = REPO_ROOT / "neuroclaw_environment.json"
            if not env_file.exists():
                await websocket.send_text(json.dumps({
                    "type": "error",
                    "message": (
                        "neuroclaw_environment.json not found. "
                        "Run python installer/setup.py to configure NeuroClaw."
                    ),
                }))
                await websocket.close()
                return

            # Initialise LLM client
            try:
                session.set_llm_client(build_llm_client(session.env))
            except Exception as exc:
                await websocket.send_text(json.dumps({
                    "type": "error",
                    "message": f"LLM backend error: {exc}",
                }))

            # Build system prompt
            soul_path = REPO_ROOT / "SOUL.md"
            soul = soul_path.read_text(encoding="utf-8") if soul_path.exists() else ""
            skill_names = ", ".join(s["name"] for s in skills)
            session.history = [
                {"role": "system", "content": f"{soul}\n\nLoaded skills: {skill_names}"}
            ]

            # Main chat loop
            try:
                while True:
                    raw = await websocket.receive_text()
                    msg = json.loads(raw)
                    user_text = msg.get("message", "").strip()
                    raw_selected = msg.get("selected_skills", [])
                    if not user_text:
                        continue

                    selected = []
                    if isinstance(raw_selected, list):
                        selected = [str(x).strip() for x in raw_selected if str(x).strip()]
                    if not selected:
                        selected = _infer_skills_from_user_text(user_text, skills)

                    selected_ctx = _selected_skills_context(selected, skills)
                    user_payload = user_text
                    if selected_ctx:
                        user_payload = (
                            f"{user_text}\n\n"
                            "[Selected skill references from local SKILL.md files]\n"
                            f"{selected_ctx}"
                        )

                    session.history.append({"role": "user", "content": user_payload})
                    try:
                        reply = await _respond(websocket, session)
                        session.history.append({"role": "assistant", "content": reply})
                    except Exception as exc:
                        err_msg = f"[Agent error: {exc}]"
                        await websocket.send_text(
                            json.dumps({"type": "error", "message": str(exc)})
                        )
                        session.history.append({"role": "assistant", "content": err_msg})

            except WebSocketDisconnect:
                pass

        except Exception as e:
            print(f"[WS] Error occurred: {type(e).__name__}: {e}", flush=True)
            import traceback
            traceback.print_exc()

    # ── Knowledge Graph Explorer ───────────────────────────────────────────

    _kg_state: dict[str, Any] = {"loaded": False, "loading": False, "error": None}
    _kg_lock = threading.Lock()

    KG_DATA_DIR = REPO_ROOT / "neurooracle" / "data"
    KG_PATH = KG_DATA_DIR / "knowledge_graph.json"
    KG_QUICK_DIR = KG_DATA_DIR / "quick"
    NEUROORACLE_HF_REPO = "zxcvb20001/NeuroOracle"
    NEUROORACLE_HF_BASE = f"https://huggingface.co/spaces/{NEUROORACLE_HF_REPO}/resolve/main"
    NEUROORACLE_MANIFEST_PATH = KG_DATA_DIR / ".download_manifest.json"
    NEUROORACLE_DOWNLOAD_FILES = [
        ("neurooracle/data/knowledge_graph.json.gz", KG_PATH.with_suffix(KG_PATH.suffix + ".gz")),
    ]
    _neurooracle_download_state: dict[str, Any] = {
        "running": False,
        "error": None,
        "completed": False,
        "current": None,
        "downloaded_bytes": 0,
        "total_bytes": None,
        "files_done": 0,
        "files_total": len(NEUROORACLE_DOWNLOAD_FILES),
    }
    _neurooracle_download_lock = threading.Lock()
    # Quick hypothesis snapshots were retired after case-study scoped runs moved
    # to explicit run directories. The KG explorer still supports hypotheses if
    # a future curated source is added here.
    HYPOTHESIS_SOURCES: list[str] = []
    RECIPES_PATH = KG_QUICK_DIR / "recipes_top10.json"

    ATOM_COLORS = {
        "disease":         "#ef4444",
        "drug":            "#ec4899",
        "imaging_marker":  "#3b82f6",
        "gene_target":     "#f59e0b",
        "cognitive_task":  "#a855f7",
        "outcome":         "#84cc16",
        "individual_data": "#14b8a6",
    }

    def _node_atoms(domain_tags) -> list[str]:
        try:
            from neurooracle.src.atoms import atoms_for_domain
        except Exception:
            return []
        seen: list[str] = []
        for d in domain_tags or []:
            for a in atoms_for_domain(d):
                if a.value not in seen:
                    seen.append(a.value)
        return seen

    RELATION_COLORS = {
        # Structural (mid + near-black slate)
        "is_a":            "#94a3b8",
        "part_of":         "#1f2937",
        # Association (sky → cyan → deep navy: max lightness spread within hue)
        "associated_with":    "#3b82f6",
        "is_associated_with": "#3b82f6",
        "correlates_with":    "#1e3a8a",
        "connects_to":        "#67e8f9",
        "projects_to":        "#155e75",
        # Causal harm (vivid red, dark wine, hot pink-red)
        "causes":          "#e63946",
        "inhibits":        "#7d1538",
        "contradicts":     "#ff4d6d",
        # Risk (orange → burnt → yellow)
        "predisposes":                  "#fb923c",
        "is_risk_factor_for":           "#c2410c",
        "gene_associated_with_disease": "#facc15",
        # Treat / mediate (emerald, very dark green, pale teal)
        "treats":          "#10b981",
        "reduces":         "#064e3b",
        "mediates":        "#5eead4",
        # Activate / boost (lime, forest, olive)
        "increases":       "#a3e635",
        "activates":       "#16a34a",
        "coactivates":     "#65a30d",
        # Modulate (royal purple)
        "modulates":       "#6d28d9",
        # Marker / predict (fuchsia, pale lavender, dark mulberry)
        "predicts":        "#d946ef",
        "is_biomarker_of": "#f0abfc",
        "distinguishes":   "#831843",
        # Provenance (very pale stone, warm dark gray)
        "about":           "#e7e5e4",
        "supported_by":    "#57534e",
    }
    DEFAULT_EDGE_COLOR = "#94a3b8"

    # ── Noise filter (query-time, does not mutate KG) ─────────────────────
    _NOISE_PREFIXES = (
        "impaired ", "increased ", "decreased ", "reduced ",
        "altered ", "elevated ", "abnormal ", "deficient ",
        "excessive ", "diminished ", "enhanced ", "disrupted ",
        "lower ", "higher ", "greater ", "lesser ",
    )
    _NOISE_SUFFIXES = (
        " findings", " levels", " changes", " symptoms",
        " deficits", " manifestations", " abnormalities",
        " dysfunctions", " status", " outcomes", " profile",
        " profiles", " patterns", " features",
    )
    NOISE_THRESHOLD = 0.3
    # Curated-vocab prefixes — trust these names even if they match noise patterns.
    # This prevents false positives on short MSH terms like "Brain", "Pons", "Sleep"
    # that trigger HypothesisEngine._is_noisy_entity's short-word regex.
    _CURATED_PREFIXES = (
        "MSH:", "NN:", "COGAT_TASK:", "COGAT_CONCEPT:", "COGAT_DISORDER:",
        "DISGENET:", "BM_REGION:", "BM_PARADIGM:", "BM_EXP:",
        "HGNC:", "NCBI_Gene:",
    )

    def _compute_noise_score(node_id: str, name: str, n_claims: int, n_hyps: int) -> float:
        """Combined heuristic score in [0, 1]. >= NOISE_THRESHOLD == noise."""
        if not name:
            return 1.0
        # Curated vocab → trust (still apply prefix/suffix check but skip token check)
        is_curated = any(node_id.startswith(p) for p in _CURATED_PREFIXES)

        score = 0.0
        if not is_curated:
            try:
                from neurooracle.src.hypothesis_engine import HypothesisEngine
                if HypothesisEngine._is_noisy_entity(name):
                    score += 0.5
            except Exception:
                pass
        lname = name.lower()
        if any(lname.startswith(p) for p in _NOISE_PREFIXES):
            score += 0.3
        if any(lname.endswith(s) for s in _NOISE_SUFFIXES):
            score += 0.3
        if node_id.startswith("CLM_CONCEPT:") and n_claims < 3:
            score += 0.15
        if n_hyps == 0 and len(name) > 40:
            score += 0.05
        return min(score, 1.0)

    def _noise_reasons(node_id: str, name: str, n_claims: int, n_hyps: int) -> list[str]:
        """Human-readable reasons why a concept was flagged as noise."""
        reasons = []
        if not name:
            return ["empty name"]
        is_curated = any(node_id.startswith(p) for p in _CURATED_PREFIXES)
        if not is_curated:
            try:
                from neurooracle.src.hypothesis_engine import HypothesisEngine
                if HypothesisEngine._is_noisy_entity(name):
                    reasons.append("generic/nominalized token (risk/effect/findings/...)")
            except Exception:
                pass
        lname = name.lower()
        for p in _NOISE_PREFIXES:
            if lname.startswith(p):
                reasons.append(f"noise prefix: '{p.strip()}'")
                break
        for s in _NOISE_SUFFIXES:
            if lname.endswith(s):
                reasons.append(f"noise suffix: '{s.strip()}'")
                break
        if node_id.startswith("CLM_CONCEPT:") and n_claims < 3:
            reasons.append("auto-extracted (CLM_CONCEPT) with <3 claims")
        if n_hyps == 0 and len(name) > 40:
            reasons.append("no hypotheses + long name")
        return reasons

    def _external_links(external_ids: dict) -> list[dict]:
        links: list[dict] = []
        for key, value in (external_ids or {}).items():
            if not value:
                continue
            url = None
            label = f"{key} · {value}"
            k = key.lower()
            v = str(value)
            if k in ("mesh_ui", "msh", "mesh"):
                url = f"https://meshb.nlm.nih.gov/record/ui?ui={v}"
                label = f"MeSH · {v}"
            elif k in ("umls", "cui"):
                url = f"https://uts.nlm.nih.gov/uts/umls/concept/{v}"
                label = f"UMLS · {v}"
            elif k in ("disgenet_id", "disgenet"):
                url = f"https://www.disgenet.org/search?source=ALL&search={v}"
                label = f"DisGeNET · {v}"
            elif k in ("nn_id", "nn", "neuronames"):
                clean = v.replace("NN:", "") if v.startswith("NN:") else v
                url = f"https://braininfo.rprc.washington.edu/centraldirectory.aspx?ID={clean}"
                label = f"NeuroNames · {clean}"
            elif k == "hgnc":
                url = f"https://www.genenames.org/data/gene-symbol-report/#!/hgnc_id/{v}"
                label = f"HGNC · {v}"
            elif k in ("ncbi_gene", "ncbi", "ncbigene"):
                url = f"https://www.ncbi.nlm.nih.gov/gene/{v}"
                label = f"NCBI Gene · {v}"
            elif k == "cogat_id":
                url = f"https://www.cognitiveatlas.org/term/id/{v}"
                label = f"Cognitive Atlas · {v}"
            elif k == "doid":
                url = f"https://disease-ontology.org/?id={v}"
                label = f"DOID · {v}"
            links.append({"key": key, "value": v, "label": label, "url": url})
        return links

    def _pmid_url(pmid: str) -> str | None:
        return f"https://pubmed.ncbi.nlm.nih.gov/{pmid}/" if pmid else None

    def _doi_url(doi: str) -> str | None:
        if not doi:
            return None
        return f"https://doi.org/{doi}" if not doi.startswith("http") else doi

    def _resolve_optional_gz(path: Path) -> Path:
        """If path doesn't exist but path.gz does, return the gz variant."""
        if path.exists():
            return path
        gz = path.with_suffix(path.suffix + ".gz")
        return gz if gz.exists() else path

    def _read_maybe_gz(path: Path) -> str:
        """Read a text file, transparently decompressing if path ends with .gz."""
        if str(path).endswith(".gz"):
            import gzip
            with gzip.open(path, "rt", encoding="utf-8") as f:
                return f.read()
        return path.read_text(encoding="utf-8")

    def _neurooracle_graph_path() -> Path:
        return _resolve_optional_gz(KG_PATH)

    def _neurooracle_proxy_url() -> str:
        return (
            os.environ.get("NEUROCLAW_PROXY_URL")
            or os.environ.get("HTTPS_PROXY")
            or os.environ.get("HTTP_PROXY")
            or os.environ.get("ALL_PROXY")
            or ""
        ).strip()

    def _neurooracle_proxy_info() -> dict[str, Any]:
        proxy_url = _neurooracle_proxy_url()
        return {
            "enabled": bool(proxy_url),
            "url": proxy_url,
            "source": "NEUROCLAW_PROXY_URL/HTTPS_PROXY/HTTP_PROXY/ALL_PROXY" if proxy_url else "",
        }

    def _neurooracle_urlopen(req: urllib.request.Request, timeout: int = 60) -> Any:
        proxy_url = _neurooracle_proxy_url()
        opener = urllib.request.build_opener(
            urllib.request.ProxyHandler({"http": proxy_url, "https": proxy_url})
        ) if proxy_url else urllib.request.build_opener()
        return opener.open(req, timeout=timeout)  # nosec: controlled NeuroOracle/HF URLs

    def _load_neurooracle_manifest() -> dict[str, Any]:
        try:
            raw = json.loads(NEUROORACLE_MANIFEST_PATH.read_text(encoding="utf-8"))
        except Exception:
            return {"files": {}}
        return raw if isinstance(raw, dict) else {"files": {}}

    def _save_neurooracle_manifest(manifest: dict[str, Any]) -> None:
        NEUROORACLE_MANIFEST_PATH.parent.mkdir(parents=True, exist_ok=True)
        NEUROORACLE_MANIFEST_PATH.write_text(
            json.dumps(manifest, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )

    def _remote_file_metadata(remote_path: str) -> dict[str, Any]:
        url = f"{NEUROORACLE_HF_BASE}/{remote_path}"
        req = urllib.request.Request(
            url,
            headers={"User-Agent": "NeuroClaw Desktop"},
            method="HEAD",
        )
        with _neurooracle_urlopen(req, timeout=20) as resp:
            headers = resp.headers
            length = headers.get("Content-Length")
            return {
                "remote_path": remote_path,
                "url": url,
                "etag": headers.get("ETag", "").strip(),
                "last_modified": headers.get("Last-Modified", "").strip(),
                "content_length": int(length) if length and length.isdigit() else None,
            }

    def _neurooracle_remote_update_status() -> dict[str, Any]:
        manifest = _load_neurooracle_manifest()
        manifest_files = manifest.get("files") if isinstance(manifest.get("files"), dict) else {}
        files: list[dict[str, Any]] = []
        update_available = False
        errors: list[str] = []

        for remote_path, dest in NEUROORACLE_DOWNLOAD_FILES:
            local_exists = dest.exists()
            local_size = dest.stat().st_size if local_exists else 0
            local_meta = manifest_files.get(remote_path, {}) if isinstance(manifest_files, dict) else {}
            item: dict[str, Any] = {
                "remote_path": remote_path,
                "local_path": str(dest),
                "local_exists": local_exists,
                "local_size": local_size,
                "update_available": False,
                "reason": "current",
            }
            try:
                remote = _remote_file_metadata(remote_path)
                item["remote"] = remote
                remote_size = remote.get("content_length")
                remote_etag = str(remote.get("etag") or "")
                remote_modified = str(remote.get("last_modified") or "")
                local_etag = str(local_meta.get("etag") or "")
                local_modified = str(local_meta.get("last_modified") or "")

                if not local_exists:
                    item.update({"update_available": True, "reason": "missing"})
                elif remote_etag and local_etag and remote_etag != local_etag:
                    item.update({"update_available": True, "reason": "etag"})
                elif remote_modified and local_modified and remote_modified != local_modified:
                    item.update({"update_available": True, "reason": "last_modified"})
                elif remote_size is not None and local_size and int(remote_size) != int(local_size):
                    item.update({"update_available": True, "reason": "size"})
                elif not local_meta:
                    item.update({"update_available": False, "reason": "no_manifest_size_match" if remote_size == local_size else "no_manifest"})
            except Exception as exc:
                item.update({"error": str(exc), "reason": "check_failed"})
                errors.append(f"{remote_path}: {exc}")

            if item.get("update_available"):
                update_available = True
            files.append(item)

        return {
            "type": "done",
            "checked_at": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime()),
            "update_available": update_available,
            "files": files,
            "errors": errors,
            "proxy": _neurooracle_proxy_info(),
            "manifest_path": str(NEUROORACLE_MANIFEST_PATH),
        }

    def _neurooracle_graph_status() -> dict[str, Any]:
        graph_path = _neurooracle_graph_path()
        exists = graph_path.exists()
        stat = graph_path.stat() if exists else None
        with _neurooracle_download_lock:
            download = dict(_neurooracle_download_state)
        return {
            "available": exists,
            "path": str(graph_path if exists else KG_PATH),
            "size_bytes": stat.st_size if stat else 0,
            "loaded": bool(_kg_state.get("loaded")),
            "loading": bool(_kg_state.get("loading")),
            "load_error": _kg_state.get("error"),
            "hf_repo": NEUROORACLE_HF_REPO,
            "proxy": _neurooracle_proxy_info(),
            "download": download,
        }

    def _reset_kg_state_after_graph_update() -> None:
        with _kg_lock:
            _kg_state.clear()
            _kg_state.update({"loaded": False, "loading": False, "error": None})

    def _download_url_to_file(url: str, dest: Path, *, label: str) -> int:
        dest.parent.mkdir(parents=True, exist_ok=True)
        tmp = dest.with_name(dest.name + ".download")
        req = urllib.request.Request(url, headers={"User-Agent": "NeuroClaw Desktop"})
        downloaded = 0
        with _neurooracle_urlopen(req, timeout=60) as resp:
            total = resp.headers.get("Content-Length")
            metadata = {
                "url": url,
                "etag": resp.headers.get("ETag", "").strip(),
                "last_modified": resp.headers.get("Last-Modified", "").strip(),
                "content_length": int(total) if total and total.isdigit() else None,
                "downloaded_at": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime()),
            }
            with _neurooracle_download_lock:
                _neurooracle_download_state["current"] = label
                _neurooracle_download_state["downloaded_bytes"] = 0
                _neurooracle_download_state["total_bytes"] = metadata["content_length"]
            with open(tmp, "wb") as f:
                while True:
                    chunk = resp.read(1024 * 1024)
                    if not chunk:
                        break
                    f.write(chunk)
                    downloaded += len(chunk)
                    with _neurooracle_download_lock:
                        _neurooracle_download_state["downloaded_bytes"] = downloaded
        os.replace(tmp, dest)
        metadata["downloaded_bytes"] = downloaded
        return metadata

    def _download_neurooracle_graph_blocking(force: bool = False) -> dict[str, Any]:
        with _neurooracle_download_lock:
            if _neurooracle_download_state.get("running"):
                already_running = True
            else:
                already_running = False
                _neurooracle_download_state.update({
                    "running": True,
                    "error": None,
                    "completed": False,
                    "current": None,
                    "downloaded_bytes": 0,
                    "total_bytes": None,
                    "files_done": 0,
                    "files_total": len(NEUROORACLE_DOWNLOAD_FILES),
                })
        if already_running:
            return _neurooracle_graph_status()

        try:
            manifest = _load_neurooracle_manifest()
            manifest_files = manifest.setdefault("files", {})
            for remote_path, dest in NEUROORACLE_DOWNLOAD_FILES:
                if dest.exists() and not force:
                    with _neurooracle_download_lock:
                        _neurooracle_download_state["files_done"] += 1
                    continue
                url = f"{NEUROORACLE_HF_BASE}/{remote_path}"
                metadata = _download_url_to_file(url, dest, label=remote_path)
                manifest_files[remote_path] = {
                    **metadata,
                    "remote_path": remote_path,
                    "local_path": str(dest),
                    "local_size": dest.stat().st_size if dest.exists() else 0,
                }
                with _neurooracle_download_lock:
                    _neurooracle_download_state["files_done"] += 1
            manifest["updated_at"] = time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())
            manifest["hf_repo"] = NEUROORACLE_HF_REPO
            _save_neurooracle_manifest(manifest)
            _reset_kg_state_after_graph_update()
            with _neurooracle_download_lock:
                _neurooracle_download_state.update({
                    "running": False,
                    "completed": True,
                    "current": None,
                    "downloaded_bytes": 0,
                    "total_bytes": None,
                })
            return _neurooracle_graph_status()
        except Exception as exc:
            with _neurooracle_download_lock:
                _neurooracle_download_state.update({
                    "running": False,
                    "completed": False,
                    "error": str(exc),
                })
            return _neurooracle_graph_status()

    def _load_kg_blocking() -> dict:
        """Load KG + hypotheses + recipes; build reverse indexes. Called once."""
        from neurooracle.src.storage import load_graph
        from neurooracle.src.hypothesis_engine import Hypothesis

        t0 = time.time()
        print(f"[kg] loading knowledge graph from {KG_PATH} ...", flush=True)
        kg = load_graph(KG_PATH)

        # name_index: lower(name|alias) -> [node_id]
        name_index: dict[str, list[str]] = {}
        concept_to_claims: dict[str, list[str]] = {}
        claim_nodes: dict[str, dict] = {}
        for nid, node in kg._index.items():
            is_claim = "claim" in node.domain_tags
            if is_claim:
                meta = node.metadata or {}
                claim_nodes[nid] = meta
                subj = meta.get("subject_id", "")
                obj = meta.get("object_id", "")
                if subj:
                    concept_to_claims.setdefault(subj, []).append(nid)
                if obj and obj != subj:
                    concept_to_claims.setdefault(obj, []).append(nid)
                continue
            # Index non-claim concepts by name and aliases
            key = node.preferred_name.strip().lower()
            if key:
                name_index.setdefault(key, []).append(nid)
            for alias in node.aliases or []:
                ak = alias.strip().lower()
                if ak and ak != key:
                    name_index.setdefault(ak, []).append(nid)

        # Load hypotheses — critic first (priority), then imaging
        hypotheses_by_id: dict[str, Hypothesis] = {}
        for fname in HYPOTHESIS_SOURCES:
            fpath = _resolve_optional_gz(KG_QUICK_DIR / fname)
            if not fpath.exists():
                # fallback: try parent data/ dir (without quick/)
                fpath = _resolve_optional_gz(KG_DATA_DIR / fname)
            if not fpath.exists():
                print(f"[kg] skip missing hypothesis file: {fname}", flush=True)
                continue
            try:
                data = json.loads(_read_maybe_gz(fpath))
                for h_dict in data.get("hypotheses", []):
                    h = Hypothesis.from_dict(h_dict)
                    if h.id and h.id not in hypotheses_by_id:
                        # tag source file for provenance
                        h.metadata = dict(h.metadata or {})
                        h.metadata.setdefault("_source_file", fname)
                        hypotheses_by_id[h.id] = h
                print(f"[kg] loaded {len(data.get('hypotheses', []))} hypotheses from {fname}", flush=True)
            except Exception as exc:
                print(f"[kg] failed to load {fname}: {exc}", flush=True)

        # Reverse index: concept_id -> [hypothesis_id]
        concept_to_hyps: dict[str, set] = {}
        for hid, h in hypotheses_by_id.items():
            touched: set[str] = set()
            if h.source_id:
                touched.add(h.source_id)
            if h.target_id:
                touched.add(h.target_id)
            for link in h.path or []:
                if link.from_id:
                    touched.add(link.from_id)
                if link.to_id:
                    touched.add(link.to_id)
            for cid in touched:
                concept_to_hyps.setdefault(cid, set()).add(hid)

        # Recipes (optional)
        recipes_by_hyp: dict[str, dict] = {}
        recipes_path = _resolve_optional_gz(RECIPES_PATH)
        if recipes_path.exists():
            try:
                rdata = json.loads(_read_maybe_gz(recipes_path))
                for r in rdata.get("recipes", []):
                    hid = r.get("hypothesis_id")
                    if hid:
                        recipes_by_hyp[hid] = r
                print(f"[kg] loaded {len(recipes_by_hyp)} recipes", flush=True)
            except Exception as exc:
                print(f"[kg] failed to load recipes: {exc}", flush=True)

        # Compute noise scores for non-claim concepts, then rebuild a clean name_index
        t_noise = time.time()
        noise_map: dict[str, float] = {}
        clean_name_index: dict[str, list[str]] = {}
        n_noisy = 0
        for nid, node in kg._index.items():
            if "claim" in (node.domain_tags or []):
                continue
            n_cl = len(concept_to_claims.get(nid, []))
            n_hy = len(concept_to_hyps.get(nid, set()))
            score = _compute_noise_score(nid, node.preferred_name or "", n_cl, n_hy)
            if score > 0:
                noise_map[nid] = score
            if score < NOISE_THRESHOLD:
                key = (node.preferred_name or "").strip().lower()
                if key:
                    clean_name_index.setdefault(key, []).append(nid)
                for alias in node.aliases or []:
                    ak = alias.strip().lower()
                    if ak and ak != key:
                        clean_name_index.setdefault(ak, []).append(nid)
            else:
                n_noisy += 1
        print(
            f"[kg] scored noise in {time.time() - t_noise:.2f}s: "
            f"{n_noisy} flagged (>= {NOISE_THRESHOLD})",
            flush=True,
        )

        stats = kg.stats()
        elapsed = time.time() - t0
        print(
            f"[kg] ready in {elapsed:.1f}s: {stats['n_concepts']} concepts, "
            f"{stats['n_edges']} edges, {len(claim_nodes)} claims, "
            f"{len(hypotheses_by_id)} hypotheses, {len(recipes_by_hyp)} recipes, "
            f"{n_noisy} noise-flagged",
            flush=True,
        )

        # ── Build trigram inverted index for fast substring search ──────
        t_tri = time.time()

        def _trigrams(s: str) -> set[str]:
            s = s.lower()
            if len(s) < 3:
                return {s} if s else set()
            return {s[i:i+3] for i in range(len(s) - 2)}

        def _build_trigram_index(idx: dict[str, list[str]]) -> dict[str, set[str]]:
            tri_idx: dict[str, set[str]] = {}
            for key in idx:
                for tri in _trigrams(key):
                    tri_idx.setdefault(tri, set()).add(key)
            return tri_idx

        trigram_index = _build_trigram_index(name_index)
        clean_trigram_index = _build_trigram_index(clean_name_index)
        print(f"[kg] built trigram indexes in {time.time() - t_tri:.2f}s", flush=True)

        # ── Pre-compute top-ranked concept lists (avoids 86k scan per request) ──
        t_top = time.time()

        def _build_top_list(idx: dict[str, list[str]], quality_strict: bool = False) -> list[dict]:
            seen: set[str] = set()
            candidates: list[dict] = []
            for nids in idx.values():
                for nid in nids:
                    if nid in seen:
                        continue
                    seen.add(nid)
                    node = kg._index.get(nid)
                    if node is None:
                        continue
                    n_cl = len(concept_to_claims.get(nid, []))
                    n_hy = len(concept_to_hyps.get(nid, set()))
                    if n_cl == 0 and n_hy == 0:
                        continue
                    if quality_strict and not (n_hy > 0 or n_cl >= 3):
                        continue
                    noise = noise_map.get(nid, 0.0)
                    candidates.append({
                        "id": nid,
                        "name": node.preferred_name,
                        "domain_tags": list(node.domain_tags or []),
                        "aliases": list(node.aliases or [])[:6],
                        "n_claims": n_cl,
                        "n_hypotheses": n_hy,
                        "noise_score": noise,
                        "is_noise": noise >= NOISE_THRESHOLD,
                    })
            candidates.sort(key=lambda r: (-(r["n_claims"] * 2 + r["n_hypotheses"]), r["noise_score"], len(r["name"])))
            return candidates

        top_all = _build_top_list(name_index)
        top_clean = _build_top_list(clean_name_index)
        print(f"[kg] pre-computed top lists in {time.time() - t_top:.2f}s ({len(top_clean)} clean, {len(top_all)} all)", flush=True)

        return {
            "loaded": True,
            "loading": False,
            "error": None,
            "kg": kg,
            "name_index": name_index,
            "clean_name_index": clean_name_index,
            "trigram_index": trigram_index,
            "clean_trigram_index": clean_trigram_index,
            "top_all": top_all,
            "top_clean": top_clean,
            "concept_to_claims": concept_to_claims,
            "claim_nodes": claim_nodes,
            "hypotheses_by_id": hypotheses_by_id,
            "concept_to_hyps": concept_to_hyps,
            "recipes_by_hyp": recipes_by_hyp,
            "noise_map": noise_map,
            "stats": {
                "n_concepts": stats["n_concepts"],
                "n_edges": stats["n_edges"],
                "n_claims": len(claim_nodes),
                "n_hypotheses": len(hypotheses_by_id),
                "n_recipes": len(recipes_by_hyp),
                "n_with_recipe": len(recipes_by_hyp),
                "n_noise_flagged": n_noisy,
                "domains": stats.get("domains", {}),
            },
        }

    async def _get_kg_state() -> dict:
        """Lazy-load the KG on first request; subsequent calls return cached state."""
        if _kg_state.get("loaded"):
            return _kg_state
        # serialize concurrent first-load attempts
        should_load = False
        with _kg_lock:
            if not _kg_state.get("loaded") and not _kg_state.get("loading"):
                _kg_state["loading"] = True
                should_load = True
        if should_load:
            try:
                new_state = await asyncio.to_thread(_load_kg_blocking)
                _kg_state.update(new_state)
            except Exception as exc:
                _kg_state["loading"] = False
                _kg_state["error"] = str(exc)
                raise
        else:
            # another request is loading — poll briefly
            for _ in range(600):  # up to ~60s
                if _kg_state.get("loaded") or _kg_state.get("error"):
                    break
                await asyncio.sleep(0.1)
        if _kg_state.get("error"):
            raise RuntimeError(_kg_state["error"])
        return _kg_state

    def _node_summary(state: dict, node_id: str) -> dict | None:
        kg = state["kg"]
        node = kg._index.get(node_id)
        if node is None:
            return None
        noise = state.get("noise_map", {}).get(node_id, 0.0)
        return {
            "id": node.id,
            "name": node.preferred_name,
            "domain_tags": list(node.domain_tags or []),
            "atoms": _node_atoms(node.domain_tags),
            "aliases": list(node.aliases or [])[:6],
            "n_claims": len(state["concept_to_claims"].get(node_id, [])),
            "n_hypotheses": len(state["concept_to_hyps"].get(node_id, set())),
            "noise_score": noise,
            "is_noise": noise >= NOISE_THRESHOLD,
        }

    def _serialize_claim(state: dict, claim_id: str) -> dict | None:
        kg = state["kg"]
        node = kg._index.get(claim_id)
        if node is None or "claim" not in node.domain_tags:
            return None
        meta = node.metadata or {}
        paper = meta.get("source_paper") or {}
        evidence = meta.get("evidence") or {}
        pmid = paper.get("pmid", "") or ""
        doi = paper.get("doi", "") or ""
        return {
            "claim_id": claim_id,
            "subject_id": meta.get("subject_id", ""),
            "subject_name": meta.get("subject_name", ""),
            "predicate": meta.get("predicate", ""),
            "object_id": meta.get("object_id", ""),
            "object_name": meta.get("object_name", ""),
            "confidence": float(meta.get("confidence", 0.0)),
            "negated": bool(meta.get("negated", False)),
            "raw_text": meta.get("raw_text", ""),
            "paper": {
                "pmid": pmid,
                "doi": doi,
                "title": paper.get("title", ""),
                "authors": paper.get("authors", ""),
                "year": paper.get("year"),
                "journal": paper.get("journal", ""),
                "pubmed_url": _pmid_url(pmid),
                "doi_url": _doi_url(doi),
            },
            "evidence": {
                "study_type": evidence.get("study_type", ""),
                "methodology": evidence.get("methodology", ""),
                "p_value": evidence.get("p_value"),
                "effect_size": evidence.get("effect_size"),
                "effect_metric": evidence.get("effect_metric", ""),
                "sample_size": evidence.get("sample_size"),
                "replicability": evidence.get("replicability", ""),
                "direction": evidence.get("direction", ""),
            },
        }

    def _serialize_hypothesis(state: dict, h, include_full_path: bool = True) -> dict:
        path_out: list[dict] = []
        pmids: set[str] = set()
        if include_full_path:
            for link in h.path or []:
                sp = link.source_paper or {}
                pmid = sp.get("pmid", "") or ""
                doi = sp.get("doi", "") or ""
                if pmid:
                    pmids.add(pmid)
                path_out.append({
                    "from_id": link.from_id,
                    "from_name": link.from_name,
                    "to_id": link.to_id,
                    "to_name": link.to_name,
                    "relation_type": link.relation_type,
                    "confidence": link.confidence,
                    "claim_id": link.claim_id,
                    "raw_text": (link.raw_text or "")[:400],
                    "paper": {
                        "pmid": pmid,
                        "doi": doi,
                        "title": sp.get("title", ""),
                        "year": sp.get("year"),
                        "journal": sp.get("journal", ""),
                        "pubmed_url": _pmid_url(pmid),
                        "doi_url": _doi_url(doi),
                    },
                })
        recipe = state["recipes_by_hyp"].get(h.id)
        return {
            "id": h.id,
            "hypothesis_type": h.hypothesis_type,
            "source_id": h.source_id,
            "source_name": h.source_name,
            "target_id": h.target_id,
            "target_name": h.target_name,
            "confidence_score": h.confidence_score,
            "novelty_score": h.novelty_score,
            "evidence_score": h.evidence_score,
            "testability_score": h.testability_score,
            "composite_score": h.composite_score,
            "critic_score": h.critic_score,
            "critic_rounds": h.critic_rounds,
            "testability_reason": h.testability_reason,
            "explanation": h.explanation,
            "path": path_out,
            "supporting_claims": list(h.supporting_claims or [])[:20],
            "source_file": (h.metadata or {}).get("_source_file", ""),
            "pmids": sorted(pmids)[:10],
            "has_recipe": recipe is not None,
            "recipe": (
                {
                    "id": recipe.get("id"),
                    "dataset": recipe.get("dataset"),
                    "model_arch": recipe.get("model_arch"),
                    "atlas": recipe.get("atlas"),
                    "target_outcome": recipe.get("target_outcome"),
                    "input_modalities": recipe.get("input_modalities"),
                    "rationale": recipe.get("rationale", "")[:300],
                }
                if recipe else None
            ),
        }

    # ── KG routes ──────────────────────────────────────────────────────────

    @app.get("/explore")
    async def kg_explore_page() -> Any:
        page = STATIC_DIR / "explore.html"
        if page.exists():
            return FileResponse(
                str(page),
                headers={
                    "Cache-Control": "no-store, no-cache, must-revalidate, max-age=0",
                    "Pragma": "no-cache",
                },
            )
        return HTMLResponse(
            "<h1>Knowledge Graph Explorer</h1>"
            f"<p>explore.html not found in <code>{STATIC_DIR}</code>.</p>",
            status_code=500,
        )

    @app.get("/api/neurooracle/graph/status")
    async def neurooracle_graph_status() -> Any:
        return _neurooracle_graph_status()

    @app.post("/api/neurooracle/graph/download")
    async def neurooracle_graph_download(payload: dict = Body(default={})) -> Any:
        force = bool((payload or {}).get("force"))
        with _neurooracle_download_lock:
            running = bool(_neurooracle_download_state.get("running"))
        if not running:
            thread = threading.Thread(
                target=_download_neurooracle_graph_blocking,
                kwargs={"force": force},
                daemon=True,
            )
            thread.start()
            # Give the worker a tiny moment to publish "running" before the UI polls.
            await asyncio.sleep(0.05)
        return _neurooracle_graph_status()

    @app.get("/api/neurooracle/graph/update-status")
    async def neurooracle_graph_update_status() -> Any:
        try:
            return await asyncio.to_thread(_neurooracle_remote_update_status)
        except Exception as exc:
            return JSONResponse(
                {"type": "error", "message": str(exc), "proxy": _neurooracle_proxy_info()},
                status_code=500,
            )

    @app.get("/api/kg/stats")
    async def kg_stats() -> Any:
        if not _neurooracle_graph_path().exists():
            return {"loaded": False, "loading": False, "missing": True, "path": str(KG_PATH)}
        if not _kg_state.get("loaded"):
            if _kg_state.get("loading"):
                return {"loaded": False, "loading": True}
            # Kick off load in background without waiting
            asyncio.create_task(_get_kg_state())
            return {"loaded": False, "loading": True}
        return {"loaded": True, **_kg_state["stats"]}

    @app.post("/api/kg/load")
    async def kg_load() -> Any:
        if not _neurooracle_graph_path().exists():
            return JSONResponse(
                {"loaded": False, "missing": True, "error": f"NeuroOracle graph file is missing: {KG_PATH}"},
                status_code=404,
            )
        try:
            state = await _get_kg_state()
            return {"loaded": True, **state["stats"]}
        except Exception as exc:
            return JSONResponse({"loaded": False, "error": str(exc)}, status_code=500)

    @app.get("/api/kg/tasks")
    async def kg_tasks() -> Any:
        try:
            from neurooracle.src.atoms import (
                CANONICAL_TASKS, CANONICAL_CHAINS, ATOM_TO_DOMAINS, Atom,
            )
        except Exception as exc:
            return JSONResponse({"error": str(exc)}, status_code=500)
        tasks = [{**t.to_dict(), "kind": "task"} for t in CANONICAL_TASKS]
        chains = [{**c.to_dict(), "kind": "chain"} for c in CANONICAL_CHAINS]
        atom_domains = {a.value: sorted(ATOM_TO_DOMAINS[a]) for a in Atom}
        return {
            "tasks": tasks,
            "chains": chains,
            "atom_domains": atom_domains,
            "atom_colors": ATOM_COLORS,
        }

    @app.get("/api/kg/task-paths")
    async def kg_task_paths(
        task: str,
        kind: str = "task",
        anchor: str = "",
        limit: int = 30,
    ) -> Any:
        """Enumerate up to `limit` atom-aligned paths through `anchor`.

        For chains, the anchor's atom pins its position in the chain (chain
        atoms are unique). We DFS forward from anchor along the suffix atoms,
        and backward along the prefix atoms, then cross-product the two halves
        into full-length paths.

        For tasks (unordered: {inputs} → output), we treat them as a 2-step
        chain anchored at query: if anchor.atom ∈ inputs, enumerate
        (anchor → output_atom_node); if anchor.atom == output, enumerate
        (input_atom_node → anchor) for each input atom. Edges may run either
        direction in the underlying KG (the relation already encodes
        directionality semantically).
        """
        try:
            from neurooracle.src.atoms import (
                CANONICAL_TASKS, CANONICAL_CHAINS, atoms_for_domain,
            )
            state = await _get_kg_state()
        except Exception as exc:
            return JSONResponse({"error": str(exc)}, status_code=500)

        kg = state["kg"]
        kind = kind.lower().strip()
        anchor_id = (anchor or "").strip()
        if not anchor_id or anchor_id not in kg._index:
            return JSONResponse({"error": "anchor required"}, status_code=400)

        # Resolve atom sequences to walk.
        if kind == "chain":
            obj = next((c for c in CANONICAL_CHAINS if c.name == task), None)
            if obj is None:
                return JSONResponse({"error": f"unknown chain: {task}"}, status_code=404)
            chain_seq = list(obj.chain)
        else:
            obj = next((t for t in CANONICAL_TASKS if t.name == task), None)
            if obj is None:
                return JSONResponse({"error": f"unknown task: {task}"}, status_code=404)
            # Build per-input atom-sequence anchored on query.
            chain_seq = None  # decided per anchor atom below

        def _atoms(nid: str) -> set:
            nd = kg._index[nid]
            out: set = set()
            for d in nd.domain_tags or []:
                out |= atoms_for_domain(d)
            return out

        anchor_atoms = _atoms(anchor_id)
        if not anchor_atoms:
            return JSONResponse({"error": f"anchor {anchor_id} has no atom"}, status_code=400)

        # Build the list of (atom_sequence) walks to enumerate. Each walk is
        # an ordered list of atoms; anchor must occupy exactly one position
        # in the sequence (its atom).
        walks: list[list] = []
        if kind == "chain":
            chain_atoms = chain_seq
            # Find anchor's index in the chain (chain atoms are unique by design).
            anchor_idx = next((i for i, a in enumerate(chain_atoms) if a in anchor_atoms), None)
            if anchor_idx is None:
                return JSONResponse(
                    {"error": f"anchor atom not in chain"}, status_code=400
                )
            walks.append((chain_atoms, anchor_idx))
        else:
            inputs = list(obj.inputs)
            output = obj.output
            if output in anchor_atoms:
                # Anchor is the output: each input becomes [input_atom, anchor_atom]
                for inp in inputs:
                    walks.append(([inp, output], 1))
            else:
                # Anchor is one of the inputs: walk anchor → output
                for inp in inputs:
                    if inp in anchor_atoms:
                        walks.append(([inp, output], 0))

        if not walks:
            return JSONResponse({"error": "no atom-walk anchored at this node"}, status_code=400)

        G = kg.G
        max_paths = max(1, min(100, int(limit)))

        # Pre-cache atoms for visited nodes to avoid recompute (atoms_for_domain is hot).
        atoms_cache: dict[str, frozenset] = {}
        def _atoms_cached(nid: str) -> frozenset:
            v = atoms_cache.get(nid)
            if v is None:
                v = frozenset(_atoms(nid))
                atoms_cache[nid] = v
            return v

        # Walk forward (next atom in sequence) along undirected adjacency, no
        # node revisits within a single path. We treat the relation as
        # undirected here: the chain's atom order is the semantic constraint;
        # KG edge direction may not always match.
        def _walk(start: str, atom_seq: list, used: set) -> list:
            """Return all simple paths whose atom-types follow atom_seq."""
            if not atom_seq:
                return [[start]]
            results: list[list] = []
            target_atom = atom_seq[0]
            neighbors = set(G.successors(start)) | set(G.predecessors(start))
            for nb in neighbors:
                if nb in used:
                    continue
                if target_atom not in _atoms_cached(nb):
                    continue
                for sub in _walk(nb, atom_seq[1:], used | {nb}):
                    results.append([start] + sub)
                    if len(results) >= max_paths * 4:  # cap explosion
                        return results
            return results

        all_paths: list[list[str]] = []
        for atom_seq, anchor_idx in walks:
            prefix_atoms = atom_seq[:anchor_idx][::-1]   # walk backward from anchor
            suffix_atoms = atom_seq[anchor_idx + 1:]     # walk forward from anchor
            prefixes = _walk(anchor_id, prefix_atoms, {anchor_id})
            suffixes = _walk(anchor_id, suffix_atoms, {anchor_id})
            for pre in prefixes:
                pre_rev = pre[::-1]  # ends at anchor
                for suf in suffixes:
                    if len(suf) > 1 and any(n in pre_rev for n in suf[1:]):
                        continue  # disjoint halves
                    full = pre_rev + suf[1:]
                    all_paths.append(full)
                    if len(all_paths) >= max_paths:
                        break
                if len(all_paths) >= max_paths:
                    break
            if len(all_paths) >= max_paths:
                break

        # De-dup paths (same node sequence).
        seen_paths: set[tuple[str, ...]] = set()
        unique_paths: list[list[str]] = []
        for p in all_paths:
            key = tuple(p)
            if key in seen_paths:
                continue
            seen_paths.add(key)
            unique_paths.append(p)
            if len(unique_paths) >= max_paths:
                break

        # Collect nodes & edges across all kept paths, with path-membership.
        node_paths: dict[str, list[int]] = {}
        edge_paths: dict[tuple[str, str], list[int]] = {}
        for pi, path in enumerate(unique_paths):
            for nid in path:
                node_paths.setdefault(nid, []).append(pi)
            for u, v in zip(path, path[1:]):
                # Use canonical (u,v) order matching DB direction if present.
                if G.has_edge(u, v):
                    key = (u, v)
                elif G.has_edge(v, u):
                    key = (v, u)
                else:
                    continue
                edge_paths.setdefault(key, []).append(pi)

        noise_map = state.get("noise_map", {})
        nodes_out: list[dict] = []
        for nid, pids in node_paths.items():
            nd = kg._index[nid]
            atoms_list = sorted(a.value for a in _atoms_cached(nid))
            primary_atom = atoms_list[0] if atoms_list else ""
            nodes_out.append({
                "id": nid,
                "label": nd.preferred_name or nid,
                "color": ATOM_COLORS.get(primary_atom, "#94a3b8"),
                "atoms": atoms_list,
                "primary_atom": primary_atom,
                "domains": list(nd.domain_tags or []),
                "size": 9,
                "noise_score": noise_map.get(nid, 0.0),
                "is_noise": noise_map.get(nid, 0.0) >= NOISE_THRESHOLD,
                "is_anchor": nid == anchor_id,
                "path_ids": pids,
            })

        edges_out: list[dict] = []
        for (s, t), pids in edge_paths.items():
            data = G.edges[s, t]
            rt = data.get("relation_type", "")
            edges_out.append({
                "id": f"e{len(edges_out)}",
                "source": s,
                "target": t,
                "label": rt,
                "color": RELATION_COLORS.get(rt, DEFAULT_EDGE_COLOR),
                "confidence": float(data.get("confidence", 1.0)),
                "path_ids": pids,
            })

        # Atom sequence for display (from the first walk; all walks share length
        # for chains; for tasks each walk is 2-atom).
        display_seq = [a.value for a in walks[0][0]] if walks else []
        return {
            "task": task,
            "kind": kind,
            "anchor": anchor_id,
            "atom_sequence": display_seq,
            "n_paths": len(unique_paths),
            "paths": [list(p) for p in unique_paths],
            "nodes": nodes_out,
            "edges": edges_out,
        }

    @app.get("/api/kg/task-subgraph")
    async def kg_task_subgraph(
        task: str,
        kind: str = "task",
        strict_chain: bool = False,
        limit: int = 200,
        anchor: str = "",
    ) -> Any:
        """Subgraph slice for a canonical task or chain.

        kind=task   → keep nodes whose atoms ∈ (inputs ∪ {output}); all edges between them.
        kind=chain  → keep nodes whose atoms ∈ chain; if strict_chain, only edges
                      whose (src_atom, tgt_atom) is an adjacent pair in chain.

        If `anchor` is given, the candidate pool is restricted to that node and
        its 1-hop neighbours (atom-filtered the same way) so the result is the
        slice of the task/chain that passes through the query node.
        """
        try:
            from neurooracle.src.atoms import (
                CANONICAL_TASKS, CANONICAL_CHAINS, atoms_for_domain,
            )
            state = await _get_kg_state()
        except Exception as exc:
            return JSONResponse({"error": str(exc)}, status_code=500)

        kg = state["kg"]
        kind = kind.lower().strip()
        if kind == "chain":
            obj = next((c for c in CANONICAL_CHAINS if c.name == task), None)
            if obj is None:
                return JSONResponse({"error": f"unknown chain: {task}"}, status_code=404)
            atom_set = set(obj.chain)
            chain_seq = list(obj.chain)
            allowed_pairs = {(chain_seq[i], chain_seq[i + 1]) for i in range(len(chain_seq) - 1)}
        else:
            obj = next((t for t in CANONICAL_TASKS if t.name == task), None)
            if obj is None:
                return JSONResponse({"error": f"unknown task: {task}"}, status_code=404)
            atom_set = set(obj.inputs) | {obj.output}
            chain_seq = []
            allowed_pairs = set()
            strict_chain = False  # not meaningful for flat task

        limit = max(20, min(800, int(limit)))

        G = kg.G
        anchor_id = (anchor or "").strip()
        if anchor_id and anchor_id not in kg._index:
            return JSONResponse({"error": f"anchor not found: {anchor_id}"}, status_code=404)

        def _node_atom_hit(nid: str) -> frozenset:
            nd = kg._index[nid]
            node_atoms: set = set()
            for d in nd.domain_tags or []:
                node_atoms |= atoms_for_domain(d)
            return frozenset(node_atoms & atom_set)

        if anchor_id:
            anchor_hit = _node_atom_hit(anchor_id)
            if not anchor_hit:
                return JSONResponse(
                    {"error": f"anchor {anchor_id} has no atom in {sorted(a.value for a in atom_set)}"},
                    status_code=400,
                )
            pool = set(G.successors(anchor_id)) | set(G.predecessors(anchor_id))
            candidates: list[tuple[int, str, frozenset]] = []
            for nid in pool:
                hit = _node_atom_hit(nid)
                if not hit:
                    continue
                candidates.append((G.degree(nid), nid, hit))
            candidates.sort(key=lambda x: x[0], reverse=True)
            kept = candidates[: max(0, limit - 1)]
            keep_ids = {nid for _, nid, _ in kept} | {anchor_id}
            node_atom_map = {nid: hit for _, nid, hit in kept}
            node_atom_map[anchor_id] = anchor_hit
            ordered_kept = [(G.degree(anchor_id), anchor_id, anchor_hit)] + kept
            n_total_candidates = len(candidates) + 1
        else:
            candidates = []
            for nid, nd in kg._index.items():
                hit = _node_atom_hit(nid)
                if not hit:
                    continue
                candidates.append((G.degree(nid), nid, hit))
            candidates.sort(key=lambda x: x[0], reverse=True)
            kept = candidates[:limit]
            keep_ids = {nid for _, nid, _ in kept}
            node_atom_map = {nid: hit for _, nid, hit in kept}
            ordered_kept = kept
            n_total_candidates = len(candidates)

        noise_map = state.get("noise_map", {})
        nodes_out: list[dict] = []
        for _, nid, hit in ordered_kept:
            nd = kg._index[nid]
            atoms_list = sorted(a.value for a in hit)
            primary_atom = atoms_list[0] if atoms_list else ""
            nodes_out.append({
                "id": nid,
                "label": nd.preferred_name or nid,
                "color": ATOM_COLORS.get(primary_atom, "#94a3b8"),
                "atoms": atoms_list,
                "primary_atom": primary_atom,
                "domains": list(nd.domain_tags or []),
                "size": 9,
                "noise_score": noise_map.get(nid, 0.0),
                "is_noise": noise_map.get(nid, 0.0) >= NOISE_THRESHOLD,
                "is_anchor": nid == anchor_id,
            })

        edges_out: list[dict] = []
        seen_pairs: set[tuple[str, str]] = set()
        for s, t, data in G.edges(data=True):
            if s not in keep_ids or t not in keep_ids:
                continue
            rt = data.get("relation_type", "")
            if not rt or rt == "about":
                continue
            if strict_chain and chain_seq:
                s_atoms = node_atom_map.get(s, frozenset())
                t_atoms = node_atom_map.get(t, frozenset())
                if not any((sa, ta) in allowed_pairs for sa in s_atoms for ta in t_atoms):
                    continue
            key = (s, t, rt)
            if key in seen_pairs:
                continue
            seen_pairs.add(key)
            edges_out.append({
                "id": f"e{len(edges_out)}",
                "source": s,
                "target": t,
                "label": rt,
                "color": RELATION_COLORS.get(rt, DEFAULT_EDGE_COLOR),
                "confidence": float(data.get("confidence", 1.0)),
            })

        return {
            "task": task,
            "kind": kind,
            "strict_chain": bool(strict_chain),
            "atoms": sorted(a.value for a in atom_set),
            "chain": [a.value for a in chain_seq] if chain_seq else [],
            "anchor": anchor_id or None,
            "nodes": nodes_out,
            "edges": edges_out,
            "truncated": n_total_candidates > len(ordered_kept),
        }

    @app.get("/api/kg/search")
    async def kg_search(
        q: str = "",
        domain: str = "",
        atom: str = "",
        limit: int = 20,
        quality: str = "clean",
    ) -> Any:
        try:
            state = await _get_kg_state()
        except Exception as exc:
            return JSONResponse({"error": str(exc)}, status_code=500)
        q_norm = (q or "").strip().lower()
        domain_filter = {d.strip() for d in (domain or "").split(",") if d.strip()}
        # Atom filter — expand each atom to its set of KG domains and union into
        # domain_filter. Unknown atom values are silently ignored.
        if atom:
            try:
                from neurooracle.src.atoms import Atom, ATOM_TO_DOMAINS
                for a in (s.strip() for s in atom.split(",") if s.strip()):
                    try:
                        domain_filter |= set(ATOM_TO_DOMAINS[Atom(a)])
                    except (ValueError, KeyError):
                        continue
            except Exception:
                pass
        quality = quality.lower() if quality else "clean"
        if quality not in ("all", "clean", "strict"):
            quality = "clean"
        idx = state["name_index"] if quality == "all" else state["clean_name_index"]
        tri_idx = state["trigram_index"] if quality == "all" else state["clean_trigram_index"]

        def passes_strict(s: dict) -> bool:
            return s["n_hypotheses"] > 0 or s["n_claims"] >= 3

        def _trigrams_q(s: str) -> set[str]:
            if len(s) < 3:
                return {s} if s else set()
            return {s[i:i+3] for i in range(len(s) - 2)}

        # Default listing when no query text: return pre-computed top list
        if len(q_norm) < 2:
            top_list = state["top_all"] if quality == "all" else state["top_clean"]
            if domain_filter:
                filtered = [r for r in top_list if domain_filter & set(r["domain_tags"])]
                if quality == "strict":
                    filtered = [r for r in filtered if passes_strict(r)]
                return {"results": filtered[:max(1, int(limit))], "query": q, "quality": quality, "mode": "top"}
            if quality == "strict":
                filtered = [r for r in top_list if passes_strict(r)]
                return {"results": filtered[:max(1, int(limit))], "query": q, "quality": quality, "mode": "top"}
            return {"results": top_list[:max(1, int(limit))], "query": q, "quality": quality, "mode": "top"}

        seen: set[str] = set()
        results: list[dict] = []

        # Exact key hit first
        for nid in idx.get(q_norm, []):
            if nid in seen:
                continue
            summary = _node_summary(state, nid)
            if summary is None:
                continue
            if domain_filter and not (domain_filter & set(summary["domain_tags"])):
                continue
            if quality == "strict" and not passes_strict(summary):
                continue
            summary["match"] = "exact"
            results.append(summary)
            seen.add(nid)

        # Trigram-accelerated substring search
        if len(results) < limit:
            tris = _trigrams_q(q_norm)
            if tris:
                candidate_keys: set[str] | None = None
                for tri in tris:
                    keys = tri_idx.get(tri)
                    if keys is None:
                        candidate_keys = set()
                        break
                    if candidate_keys is None:
                        candidate_keys = set(keys)
                    else:
                        candidate_keys &= keys
                for key in (candidate_keys or set()):
                    if q_norm not in key:
                        continue
                    for nid in idx.get(key, []):
                        if nid in seen:
                            continue
                        summary = _node_summary(state, nid)
                        if summary is None:
                            continue
                        if domain_filter and not (domain_filter & set(summary["domain_tags"])):
                            continue
                        if quality == "strict" and not passes_strict(summary):
                            continue
                        summary["match"] = "substring"
                        results.append(summary)
                        seen.add(nid)
                        if len(results) >= limit * 3:
                            break
                    if len(results) >= limit * 3:
                        break

        results.sort(
            key=lambda r: (
                0 if r.get("match") == "exact" else 1,
                r.get("noise_score", 0.0),
                -(r["n_hypotheses"] * 2 + r["n_claims"]),
                len(r["name"]),
            )
        )
        return {"results": results[:limit], "query": q, "quality": quality, "mode": "search"}

    @app.get("/api/kg/node/{node_id}")
    async def kg_node(node_id: str) -> Any:
        try:
            state = await _get_kg_state()
        except Exception as exc:
            return JSONResponse({"error": str(exc)}, status_code=500)
        kg = state["kg"]
        node = kg._index.get(node_id)
        if node is None:
            return JSONResponse({"error": f"node not found: {node_id}"}, status_code=404)
        n_claims = len(state["concept_to_claims"].get(node_id, []))
        n_hyps = len(state["concept_to_hyps"].get(node_id, set()))
        noise = state.get("noise_map", {}).get(node_id, 0.0)
        reasons = _noise_reasons(node_id, node.preferred_name or "", n_claims, n_hyps) if noise >= NOISE_THRESHOLD else []
        return {
            "id": node.id,
            "name": node.preferred_name,
            "definition": node.definition or "",
            "domain_tags": list(node.domain_tags or []),
            "atoms": _node_atoms(node.domain_tags),
            "semantic_types": list(node.semantic_types or []),
            "source_vocab": node.source_vocab or "",
            "aliases": list(node.aliases or []),
            "external_ids": dict(node.external_ids or {}),
            "external_links": _external_links(node.external_ids),
            "atlas_mapping": node.atlas_mapping,
            "n_claims": n_claims,
            "n_hypotheses": n_hyps,
            "noise_score": noise,
            "is_noise": noise >= NOISE_THRESHOLD,
            "noise_reasons": reasons,
            "color": ATOM_COLORS.get(
                (_node_atoms(node.domain_tags) or [""])[0], "#94a3b8"
            ),
        }

    @app.get("/api/kg/node/{node_id}/neighborhood")
    async def kg_neighborhood(
        node_id: str,
        depth: int = 1,
        edge_types: str = "",
        limit: int = 80,
    ) -> Any:
        try:
            state = await _get_kg_state()
        except Exception as exc:
            return JSONResponse({"error": str(exc)}, status_code=500)
        kg = state["kg"]
        if node_id not in kg._index:
            return JSONResponse({"error": f"node not found: {node_id}"}, status_code=404)

        depth = max(1, min(2, int(depth)))
        limit = max(10, min(200, int(limit)))
        type_filter = {t.strip() for t in edge_types.split(",") if t.strip() and t.strip() != "all"}

        G = kg.G
        visited = {node_id}
        depth_map: dict[str, int] = {node_id: 0}
        frontier = {node_id}
        edges_collected: list[tuple[str, str, dict]] = []
        for hop in range(depth):
            next_frontier: set[str] = set()
            for n in frontier:
                for _, tgt, data in G.out_edges(n, data=True):
                    rt = data.get("relation_type", "")
                    if type_filter and rt not in type_filter:
                        continue
                    if rt == "about":
                        continue
                    if hop > 0 and tgt in visited:
                        continue
                    edges_collected.append((n, tgt, data))
                    if tgt not in visited:
                        next_frontier.add(tgt)
                        depth_map.setdefault(tgt, hop + 1)
                for src, _, data in G.in_edges(n, data=True):
                    rt = data.get("relation_type", "")
                    if type_filter and rt not in type_filter:
                        continue
                    if rt == "about":
                        continue
                    if hop > 0 and src in visited:
                        continue
                    edges_collected.append((src, n, data))
                    if src not in visited:
                        next_frontier.add(src)
                        depth_map.setdefault(src, hop + 1)
            visited |= next_frontier
            # For depth>=2: after hop 0, narrow the frontier to the top-N
            # depth-1 nodes by degree so hop 1 doesn't explode, but still runs.
            if depth >= 2 and hop == 0 and len(next_frontier) > limit:
                ranked = sorted(next_frontier, key=lambda n: G.degree(n), reverse=True)
                frontier = set(ranked[:limit])
            else:
                frontier = next_frontier

        # depth=2: drop edges that connect two nodes already at depth=1 to avoid
        # cluttering the graph with sibling cross-links (the user asked that
        # depth=2 not consider two peer neighbors being connected). Keep edges
        # that touch the center or a depth=2 node.
        if depth >= 2:
            edges_collected = [
                (s, t, d) for (s, t, d) in edges_collected
                if not (depth_map.get(s) == 1 and depth_map.get(t) == 1)
            ]

        # Rank candidate nodes by degree, but when depth>=2 reserve slots for
        # depth-2 nodes so the user actually sees A→B→C chains.
        node_ids: list[str] = [node_id]
        slots = max(0, limit - 1)
        if depth >= 2:
            d1 = [n for n in visited if n != node_id and depth_map.get(n) == 1]
            d2 = [n for n in visited if n != node_id and depth_map.get(n) == 2]
            d1.sort(key=lambda n: G.degree(n), reverse=True)
            d2.sort(key=lambda n: G.degree(n), reverse=True)
            # Reserve ~40% of slots for depth-2 nodes (at least 5 if available)
            d2_slots = max(5, slots * 2 // 5)
            d1_slots = slots - min(d2_slots, len(d2))
            node_ids.extend(d1[:d1_slots])
            node_ids.extend(d2[:slots - len(node_ids) + 1])
        else:
            candidates = [n for n in visited if n != node_id]
            candidates.sort(key=lambda n: G.degree(n), reverse=True)
            node_ids.extend(candidates[:slots])
        keep = set(node_ids)

        nodes_out: list[dict] = []
        noise_map = state.get("noise_map", {})
        for nid in node_ids:
            nd = kg._index.get(nid)
            if nd is None:
                continue
            is_claim = "claim" in (nd.domain_tags or [])
            atoms_n = _node_atoms(nd.domain_tags)
            primary_atom = atoms_n[0] if atoms_n else ""
            label = nd.preferred_name or nid
            if is_claim and len(label) > 60:
                label = label[:57] + "…"
            noise = noise_map.get(nid, 0.0)
            is_noisy = noise >= NOISE_THRESHOLD
            base_size = 14 if nid == node_id else (6 if is_claim else 9)
            nodes_out.append({
                "id": nid,
                "label": label,
                "color": ATOM_COLORS.get(primary_atom, "#94a3b8"),
                "domain": primary_atom,
                "domains": list(nd.domain_tags or []),
                "atoms": atoms_n,
                "is_claim": is_claim,
                "is_center": nid == node_id,
                "depth": depth_map.get(nid, 0),
                "size": base_size if not (is_noisy and nid != node_id) else max(3, int(base_size * 0.55)),
                "noise_score": noise,
                "is_noise": is_noisy,
            })

        # Aggregate edges by unordered pair so bidirectional or multi-predicate
        # edges render as a single visual line (prevents label overlap).
        # Additionally scan claim nodes to surface predicates that the DiGraph
        # collapsed (graph_manager keeps only the highest-confidence relation).
        pair_info: dict[frozenset, dict] = {}
        for src, tgt, data in edges_collected:
            if src not in keep or tgt not in keep:
                continue
            rt = data.get("relation_type", "")
            if not rt:
                continue
            pair = frozenset({src, tgt}) if src != tgt else frozenset({src})
            entry = pair_info.setdefault(pair, {
                "src": src, "tgt": tgt,   # may be overwritten; used for first-seen direction
                "relations_fwd": [],     # ordered, deduped
                "relations_rev": [],
                "confidence": 0.0,
            })
            entry["confidence"] = max(entry["confidence"], float(data.get("confidence", 1.0)))
            # Track which direction this relation was seen in relative to (src, tgt)
            if (src, tgt) == (entry["src"], entry["tgt"]):
                if rt not in entry["relations_fwd"]:
                    entry["relations_fwd"].append(rt)
            else:
                if rt not in entry["relations_rev"]:
                    entry["relations_rev"].append(rt)

        # Pull additional claim-backed predicates between kept pairs
        concept_to_claims = state["concept_to_claims"]
        for pair, entry in pair_info.items():
            a_raw = list(pair)
            if len(a_raw) == 1:
                continue  # self-loop; skip extra claim scan
            a, b = a_raw[0], a_raw[1]
            a_claims = set(concept_to_claims.get(a, []))
            b_claims = set(concept_to_claims.get(b, []))
            shared = a_claims & b_claims
            if not shared:
                continue
            for cid in shared:
                cn = kg._index.get(cid)
                if cn is None:
                    continue
                meta = cn.metadata or {}
                pred = (meta.get("predicate") or "").strip()
                if not pred:
                    continue
                if type_filter and pred not in type_filter:
                    continue
                subj = meta.get("subject_id", "")
                obj = meta.get("object_id", "")
                if (subj, obj) == (entry["src"], entry["tgt"]):
                    if pred not in entry["relations_fwd"]:
                        entry["relations_fwd"].append(pred)
                elif (subj, obj) == (entry["tgt"], entry["src"]):
                    if pred not in entry["relations_rev"]:
                        entry["relations_rev"].append(pred)

        # Emit merged edges
        edges_out: list[dict] = []
        for pair, entry in pair_info.items():
            fwd = entry["relations_fwd"]
            rev = entry["relations_rev"]
            if not fwd and not rev:
                continue
            # Combine labels. If bidirectional, join both with ⇄ so the user sees
            # there are multiple relations.
            parts: list[str] = []
            if fwd:
                parts.append(" · ".join(fwd[:3]) + (f" +{len(fwd)-3}" if len(fwd) > 3 else ""))
            if rev:
                parts.append("← " + " · ".join(rev[:3]) + (f" +{len(rev)-3}" if len(rev) > 3 else ""))
            label = "  ⇄  ".join(parts) if (fwd and rev) else (parts[0] if parts else "")
            primary = (fwd[0] if fwd else rev[0])
            edges_out.append({
                "id": f"e{len(edges_out)}",
                "source": entry["src"],
                "target": entry["tgt"],
                "label": label,
                "relations_fwd": fwd,
                "relations_rev": rev,
                "bidirectional": bool(fwd and rev),
                "color": RELATION_COLORS.get(primary, DEFAULT_EDGE_COLOR),
                "confidence": entry["confidence"],
            })

        return {
            "center": node_id,
            "depth": depth,
            "nodes": nodes_out,
            "edges": edges_out,
            "depth_map": {nid: depth_map.get(nid, 0) for nid in keep},
            "truncated": len(visited) > limit,
        }

    @app.get("/api/kg/node/{node_id}/claims")
    async def kg_claims(
        node_id: str,
        limit: int = 50,
        predicate: str = "",
        neighbor_id: str = "",
    ) -> Any:
        try:
            state = await _get_kg_state()
        except Exception as exc:
            return JSONResponse({"error": str(exc)}, status_code=500)
        kg = state["kg"]
        if node_id not in kg._index:
            return JSONResponse({"error": f"node not found: {node_id}"}, status_code=404)
        claim_ids = state["concept_to_claims"].get(node_id, [])
        # parse predicate filter (may be comma-separated list from the UI)
        pred_filter = {p.strip() for p in (predicate or "").split(",") if p.strip() and p.strip() != "all"}
        items = []
        for cid in claim_ids:
            s = _serialize_claim(state, cid)
            if not s:
                continue
            if pred_filter and s.get("predicate") not in pred_filter:
                continue
            if neighbor_id:
                if s.get("subject_id") != neighbor_id and s.get("object_id") != neighbor_id:
                    continue
            items.append(s)
        # Sort: confidence desc, year desc
        items.sort(key=lambda c: (
            -(c.get("confidence") or 0.0),
            -((c.get("paper") or {}).get("year") or 0),
        ))
        return {"node_id": node_id, "total": len(items), "claims": items[: max(1, int(limit))]}

    @app.get("/api/kg/edge-sources")
    async def kg_edge_sources(source: str = "", target: str = "", limit: int = 50) -> Any:
        """Return all claims + curated edges that connect two concepts (either direction)."""
        try:
            state = await _get_kg_state()
        except Exception as exc:
            return JSONResponse({"error": str(exc)}, status_code=500)
        kg = state["kg"]
        if not source or not target:
            return JSONResponse({"error": "source and target are required"}, status_code=400)
        if source not in kg._index or target not in kg._index:
            return JSONResponse({"error": "node(s) not found"}, status_code=404)

        src_node = kg._index[source]
        tgt_node = kg._index[target]

        # 1. Claims where {subject, object} == {source, target}
        claim_items: list[dict] = []
        seen_claims: set[str] = set()
        src_claims = set(state["concept_to_claims"].get(source, []))
        tgt_claims = set(state["concept_to_claims"].get(target, []))
        for cid in src_claims & tgt_claims:
            if cid in seen_claims:
                continue
            seen_claims.add(cid)
            s = _serialize_claim(state, cid)
            if s:
                claim_items.append(s)
        claim_items.sort(key=lambda c: (
            -(c.get("confidence") or 0.0),
            -((c.get("paper") or {}).get("year") or 0),
        ))

        # 2. Curated edges (non-claim) between the two nodes, both directions
        curated_edges: list[dict] = []
        G = kg.G
        for u, v in ((source, target), (target, source)):
            if G.has_edge(u, v):
                data = G.edges[u, v]
                src_str = data.get("source", "")
                if src_str.startswith("claim:"):
                    continue  # already counted above
                curated_edges.append({
                    "from_id": u,
                    "from_name": kg._index[u].preferred_name,
                    "to_id": v,
                    "to_name": kg._index[v].preferred_name,
                    "relation_type": data.get("relation_type", ""),
                    "confidence": float(data.get("confidence", 1.0)),
                    "source_vocab": src_str or "curated",
                    "evidence_ref": data.get("evidence_ref", ""),
                })

        return {
            "source": {"id": source, "name": src_node.preferred_name},
            "target": {"id": target, "name": tgt_node.preferred_name},
            "total_claims": len(claim_items),
            "total_curated_edges": len(curated_edges),
            "claims": claim_items[: max(1, int(limit))],
            "curated_edges": curated_edges,
        }

    @app.post("/api/kg/path-claims")
    async def kg_path_claims(payload: dict = Body(...)) -> Any:
        """Given a list of node-paths, return claims grouped per adjacent pair.

        Body: {"paths": [["A","B","C"], ["A","D","C"], ...], "limit": 30}
        Response: {"edges": [{"from_id","from_name","to_id","to_name","claims":[...],"curated_edges":[...]}]}
        Pairs are deduplicated across paths (each unique unordered pair shows once).
        """
        try:
            state = await _get_kg_state()
        except Exception as exc:
            return JSONResponse({"error": str(exc)}, status_code=500)
        kg = state["kg"]
        paths = payload.get("paths") or []
        per_edge_limit = max(1, min(100, int(payload.get("limit") or 30)))

        seen_pairs: set[tuple[str, str]] = set()
        edges_out: list[dict] = []
        for path in paths:
            if not isinstance(path, list):
                continue
            for u, v in zip(path, path[1:]):
                if not isinstance(u, str) or not isinstance(v, str):
                    continue
                key = (u, v) if u < v else (v, u)
                if key in seen_pairs:
                    continue
                seen_pairs.add(key)
                if u not in kg._index or v not in kg._index:
                    continue

                # Same logic as /api/kg/edge-sources but inlined to batch.
                u_claims = set(state["concept_to_claims"].get(u, []))
                v_claims = set(state["concept_to_claims"].get(v, []))
                claim_items: list[dict] = []
                for cid in u_claims & v_claims:
                    s = _serialize_claim(state, cid)
                    if s:
                        claim_items.append(s)
                claim_items.sort(key=lambda c: (
                    -(c.get("confidence") or 0.0),
                    -((c.get("paper") or {}).get("year") or 0),
                ))

                curated_edges: list[dict] = []
                G = kg.G
                for a, b in ((u, v), (v, u)):
                    if G.has_edge(a, b):
                        data = G.edges[a, b]
                        src_str = data.get("source", "")
                        if src_str.startswith("claim:"):
                            continue
                        curated_edges.append({
                            "from_id": a, "from_name": kg._index[a].preferred_name,
                            "to_id": b, "to_name": kg._index[b].preferred_name,
                            "relation_type": data.get("relation_type", ""),
                            "confidence": float(data.get("confidence", 1.0)),
                            "source_vocab": src_str or "curated",
                            "evidence_ref": data.get("evidence_ref", ""),
                        })

                edges_out.append({
                    "from_id": u, "from_name": kg._index[u].preferred_name,
                    "to_id": v, "to_name": kg._index[v].preferred_name,
                    "total_claims": len(claim_items),
                    "claims": claim_items[:per_edge_limit],
                    "curated_edges": curated_edges,
                })

        return {"edges": edges_out, "n_edges": len(edges_out)}

    @app.get("/api/kg/node/{node_id}/hypotheses")
    async def kg_hypotheses(
        node_id: str,
        limit: int = 20,
        min_score: float = 0.0,
        recipe_only: bool = False,
    ) -> Any:
        try:
            state = await _get_kg_state()
        except Exception as exc:
            return JSONResponse({"error": str(exc)}, status_code=500)
        kg = state["kg"]
        if node_id not in kg._index:
            return JSONResponse({"error": f"node not found: {node_id}"}, status_code=404)
        hyp_ids = state["concept_to_hyps"].get(node_id, set())
        hyps = [state["hypotheses_by_id"][hid] for hid in hyp_ids if hid in state["hypotheses_by_id"]]
        # Filters
        if min_score > 0:
            hyps = [h for h in hyps if (h.composite_score or 0.0) >= min_score]
        if recipe_only:
            hyps = [h for h in hyps if h.id in state["recipes_by_hyp"]]
        hyps.sort(key=lambda h: (h.composite_score or 0.0), reverse=True)
        items = [_serialize_hypothesis(state, h) for h in hyps[: max(1, int(limit))]]
        return {
            "node_id": node_id,
            "total": len(hyps),
            "hypotheses": items,
            "has_recipes": len(state["recipes_by_hyp"]) > 0,
        }

    @app.get("/api/kg/hypothesis/{hyp_id}")
    async def kg_hypothesis_detail(hyp_id: str) -> Any:
        try:
            state = await _get_kg_state()
        except Exception as exc:
            return JSONResponse({"error": str(exc)}, status_code=500)
        h = state["hypotheses_by_id"].get(hyp_id)
        if h is None:
            return JSONResponse({"error": f"hypothesis not found: {hyp_id}"}, status_code=404)
        return _serialize_hypothesis(state, h)

    @app.get("/api/kg/paths")
    async def kg_paths(
        source: str,
        target: str,
        max_hops: int = 3,
        max_paths: int = 50,
        directed: bool = False,
        exclude_predicates: str = "about,supported_by,contradicts",
        skip_high_degree: int = 2000,
        timeout_ms: int = 4000,
        count_cap: int = 2000,
    ) -> Any:
        """Find simple paths between two concept nodes.

        On-demand only; not auto-fired. Bounded by max_hops / max_paths /
        timeout to keep responses fast even on hubs like 'brain' or 'fmri'.

        Returns up to `max_paths` serialized paths plus `total_paths`
        (the full count discovered up to `count_cap`, even paths beyond
        the displayed window — gives the user a sense of how many exist).
        """
        try:
            state = await _get_kg_state()
        except Exception as exc:
            return JSONResponse({"error": str(exc)}, status_code=500)
        kg = state["kg"]
        if source not in kg._index:
            return JSONResponse({"error": f"source not found: {source}"}, status_code=404)
        if target not in kg._index:
            return JSONResponse({"error": f"target not found: {target}"}, status_code=404)
        if source == target:
            return JSONResponse({"error": "source and target are identical"}, status_code=400)

        max_hops = max(1, min(5, int(max_hops)))
        max_paths = max(1, min(500, int(max_paths)))
        count_cap = max(max_paths, min(20000, int(count_cap)))
        timeout_s = max(0.5, min(10.0, float(timeout_ms) / 1000.0))
        skip_high_degree = max(50, int(skip_high_degree))
        excl = {p.strip() for p in exclude_predicates.split(",") if p.strip()}

        G = kg.G
        index = kg._index

        def _dfs() -> dict:
            """Synchronous DFS; runs in a worker thread."""
            import time as _time

            # `paths` stores at most max_paths serialized; `total_count` keeps
            # counting beyond that so we can show "displaying X of Y total".
            paths: list[tuple[list[str], list[tuple[str, str, dict, str]]]] = []
            total_count = 0
            deadline = _time.monotonic() + timeout_s
            truncated_by = None  # "count_cap" | "timeout" | None
            visited_count = 0

            def neighbors(node: str):
                for _, v, edata in G.out_edges(node, data=True):
                    yield v, edata, "out"
                if not directed:
                    for u, _, edata in G.in_edges(node, data=True):
                        yield u, edata, "in"

            def expand(
                node: str,
                path_nodes: list[str],
                path_edges: list[tuple[str, str, dict, str]],
                depth: int,
                on_path: set[str],
            ) -> None:
                nonlocal truncated_by, visited_count, total_count
                if truncated_by is not None:
                    return
                visited_count += 1
                if visited_count & 0xFF == 0 and _time.monotonic() > deadline:
                    truncated_by = "timeout"
                    return
                if node == target and depth > 0:
                    total_count += 1
                    if len(paths) < max_paths:
                        paths.append((list(path_nodes), list(path_edges)))
                    if total_count >= count_cap:
                        truncated_by = "count_cap"
                    return
                if depth >= max_hops:
                    return
                for v, edata, direction in neighbors(node):
                    rt = edata.get("relation_type", "")
                    if rt in excl:
                        continue
                    if v in on_path:
                        continue
                    nv = index.get(v)
                    if nv is not None and "claim" in (nv.domain_tags or []):
                        continue
                    if v != target and v != source:
                        try:
                            if G.degree(v) > skip_high_degree:
                                continue
                        except Exception:
                            pass
                    on_path.add(v)
                    path_nodes.append(v)
                    path_edges.append((node, v, edata, direction))
                    expand(v, path_nodes, path_edges, depth + 1, on_path)
                    on_path.discard(v)
                    path_nodes.pop()
                    path_edges.pop()

            expand(source, [source], [], 0, {source})

            serialized = []
            for pn, pe in paths:
                nodes_out = []
                for nid in pn:
                    n = index.get(nid)
                    if n is None:
                        nodes_out.append({"id": nid, "name": nid, "atom": ""})
                    else:
                        atoms_n = _node_atoms(n.domain_tags)
                        nodes_out.append({
                            "id": nid,
                            "name": n.preferred_name,
                            "atom": atoms_n[0] if atoms_n else "",
                        })
                edges_out = []
                conf_sum = 0.0
                for s, t, edata, direction in pe:
                    conf = float(edata.get("confidence") or 0.0)
                    conf_sum += conf
                    edges_out.append({
                        "source": s,
                        "target": t,
                        "predicate": edata.get("relation_type", ""),
                        "direction": direction,
                        "confidence": conf,
                        "edge_source": edata.get("source", ""),
                    })
                serialized.append({
                    "nodes": nodes_out,
                    "edges": edges_out,
                    "length": len(edges_out),
                    "avg_confidence": (conf_sum / len(edges_out)) if edges_out else 0.0,
                })

            serialized.sort(key=lambda p: (p["length"], -p["avg_confidence"]))

            return {
                "paths": serialized,
                "total_count": total_count,
                "truncated_by": truncated_by,
                "expansions": visited_count,
            }

        result = await asyncio.to_thread(_dfs)
        return {
            "source": source,
            "target": target,
            "max_hops": max_hops,
            "directed": directed,
            "displayed_paths": len(result["paths"]),
            "total_paths": result["total_count"],
            "paths": result["paths"],
            "truncated_by": result["truncated_by"],
            "expansions": result["expansions"],
        }

    return app


# ── Entry point ────────────────────────────────────────────────────────────────

def run_server(host: str = DEFAULT_HOST, port: int = DEFAULT_PORT) -> None:
    """Start the uvicorn server (blocking call — returns only when the server stops)."""
    _require_webdeps()
    import uvicorn  # type: ignore

    app = create_app()
    print(f"\n  NeuroClaw Web UI  →  http://{host}:{port}\n")
    uvicorn.run(app, host=host, port=port, log_level="info")


def main() -> None:
    parser = argparse.ArgumentParser(
        description="NeuroClaw Web UI — start the browser-based chat interface."
    )
    parser.add_argument(
        "--host", default=DEFAULT_HOST,
        help=f"Bind host (default: {DEFAULT_HOST})",
    )
    parser.add_argument(
        "--port", type=int, default=DEFAULT_PORT,
        help=f"Port number (default: {DEFAULT_PORT})",
    )
    args = parser.parse_args()
    run_server(host=args.host, port=args.port)


if __name__ == "__main__":
    main()
